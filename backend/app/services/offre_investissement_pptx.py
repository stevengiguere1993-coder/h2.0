"""Génère un .pptx d'offre d'investissement Horizon depuis une
`LeadAnalysis` + un input humain de stratégie value-add.

Approche :
    * Le template `horizon_v2.pptx` (par défaut, configurable via env var
      `OFFRE_INVEST_TEMPLATE`) est une copie du deck 1660 Saint-Clément
      version updated (16 slides, avec slide TENDANCES en pos. 11).
      Les valeurs littérales (montants, phrases, %, ...) servent de
      "placeholders" identifiables — une chaîne dans le template = une
      variable du mapping `horizon_v2.mapping.json`.
    * On effectue une substitution `string.replace` au niveau des runs et
      cellules de tableau pour préserver polices, couleurs, charts, layout.
    * Les variables sont catégorisées en 3 sources :
        - `auto`     : extrait direct de `LeadAnalysis` (~30 champs)
        - `hybrid`   : recalculé depuis le moteur d'analyse financière +
                       inputs value-add (loyers projetés, AT/RT, etc.)
        - `human`    : saisi par l'utilisateur dans le wizard (tagline,
                       bullets opportunité, libellés, qualificatif, etc.)
    * Les photos remplacent les blobs existants en gardant le rectangle
      (4 slots : cover, exterieur, carte, tendances).
    * Charts embedded (slide 4 / Opportunité unique, slide 12 / Tendances,
      slide 13 / Valeur ajoutée investisseur) : substitution des séries
      via `chart.replace_data(CategoryChartData)` — préserve les couleurs
      et le formatting visuel du template.
    * Logique conditionnelle :
        - Équité refi négative → libellé « Remboursement du partenaire »
          + pourcentage au lieu de montant $
        - Taxes scolaires == 0 → libellé "(inclus scolaire)" préservé

Versioning :
    * v1 (horizon_v1.pptx) : conservé pour rétrocompatibilité MVP PR #532
    * v2 (horizon_v2.pptx) : 16 slides + TENDANCES + 10 vars (PR #533)
    * v3 (logique service, même template `horizon_v2.pptx`) : corrige
      slides 3/4/5/6/9/10/11/12 (PR #534)
    * v4 (logique service + template `horizon_v2.pptx` modifié) :
      perfectionne slides 3-6 (PR de ce travail) :
        - slide 3 (idx 2) :
            * prix d'acquisition forcé sur asking_price via accès direct
              à la cellule Table 11[0][1] (source documentée et garantie)
            * nb d'unités préserve le formatage taille 14 pas-gras via
              accès direct à la cellule Table 12[0][1] (multi-paragraphes)
        - slide 4 (idx 3) :
            * bullet 3 (paragraphe 2 du TextBox 30) :
                - remplace SEULEMENT le contenu via accès direct au
                  paragraphe, sans préfixe "Levier principal:"
                - format forcé : bold=False, color=blanc (FFFFFF),
                  size=14pt pour matcher les bullets 1/2/4
            * Chart 79 (VALEUR MARCHANDE LINE) :
                - "Avant" = asking_price (= prix payé)
                - "Après refi" = scenario_choisi.valeur_economique
                  (APH 50 ou APH 100 selon programme_schl coché)
                - échelle Y auto (template nettoyé du min hardcodé)
            * Chart 85 (PROFIT À L'ACHAT COLUMN STACKED) :
                - Catégorie 2 renommée "Valeur réelle" (au lieu de
                  "Valeur municipale")
                - Valeur 2 = valeur_comparable_centris (nouveau champ
                  wizard) ou fallback evaluation_municipale
                - 2 séries empilées : "Prix payé" (jaune F8D956) +
                  "Écart comparable" (rouge EF3E45)
                - échelle Y auto
        - slide 5 (idx 4) :
            * sources documentées par commentaires inline pour chaque
              cellule des tableaux Avant/Après (cohérence avec slides
              3 et 4)
            * coût énergétique "Avant" = energie de la fiche
            * coût énergétique "Après" = 0 si conversion_chauffage
              elec_to_thermo, sinon energie
        - slide 6 (idx 5) :
            * "Création de chambre" affichée conditionnellement
              UNIQUEMENT si conversion_chambres=True :
                - Table 1 cell[4][1] "Création chambres" → label
                  selon stratégie ou vide
                - Table 12 cell[1][1] "Fin de création des chambres" →
                  libellé alternatif ou vide selon stratégie
    * v5a (logique service uniquement, template inchangé) : corrige
      slides 8, 9, 10 (visibles = idx 7, 8, 9 dans prs.slides) :
        - slide 8 (idx 7) FINANCES ACHAT :
            * Table 2 cell[4][1] "Nombre de logements" : substitué
              par `rec.nb_logements` (avant hardcodé à "8")
        - slide 9 (idx 8) FINANCES OPTIMISATION :
            * Frais autres (col 7 rows 5-11) substitués cellule par
              cellule : évaluateur / évaluateur 2 / inspection /
              avocat / notaire / notaire 2 / rapport efficacité.
              Source : `rec.frais_demarrage_overrides_json` (par fiche)
              avec fallback sur les défauts industrie (2000$/2000$/
              3000$/5000$/2000$/2000$/5000$).
            * Total frais autres cell[13][7] recalculé avec les
              valeurs effectivement utilisées (au lieu du « 108 000$ »
              du template 1660).
            * Tableau rénovations condensé : les 16 slots rénos
              (rows 2..17) dont le libellé est vide sont retirés
              physiquement APRÈS substitution des placeholders →
              le Total reste collé aux items utilisés.
        - slide 10 (idx 9) FINANCES REFINANCEMENT :
            * Table 2 cell[5][1] "Nombre de logements" : substitué
              par `rec.nb_logements` (avant hardcodé à "8")
            * Format `_money_long_signed` unifié sans espace de tête
              (-65 000$ au lieu de ` -65 000$`)
            * Cellules équité [9][4] et [9][7] : si l'équité du
              scénario est NÉGATIVE, on affiche un pourcentage
              cohérent avec le callout TextBox 18 (« Jusqu'à XX% »)
              au lieu d'un montant négatif. Si positive, montant en $.
    * Env var `OFFRE_INVEST_TEMPLATE` (défaut: 'horizon_v2') permet de
      switcher.

Limites :
    * Les recalculs hybrides utilisent les RÉSULTATS DÉJÀ PERSISTÉS dans
      `analysis_results_json` (le run-financial-analysis le plus récent)
      + des ajustements simples pour les value-add (nouveau loyer moyen,
      revenus, RCI/PVI). Pas de recalcul temps-réel du moteur SCHL/APH.
"""

from __future__ import annotations

import io
import json
import logging
import os
import re
import unicodedata
from dataclasses import dataclass, field
from datetime import date, datetime, timedelta
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from sqlalchemy.ext.asyncio import AsyncSession

from app.models.lead_analysis import LeadAnalysis, LeadAnalysisAttachment


log = logging.getLogger(__name__)


_TEMPLATE_DIR = (
    Path(__file__).resolve().parent.parent
    / "templates"
    / "offre_investissement"
)


def _resolve_template_paths() -> Tuple[Path, Path, str]:
    """Résout le couple (template.pptx, mapping.json, version) selon
    l'env var ``OFFRE_INVEST_TEMPLATE`` (défaut: ``horizon_v2``).

    Fallback : si le template demandé n'existe pas sur disque, on
    retombe sur ``horizon_v2`` puis ``horizon_v1`` pour ne jamais
    bloquer la génération (utile en cas de déploiement partiel).
    """
    version = os.environ.get("OFFRE_INVEST_TEMPLATE", "horizon_v2").strip()
    if not version:
        version = "horizon_v2"
    candidates = [version, "horizon_v2", "horizon_v1"]
    seen: set = set()
    for v in candidates:
        if v in seen:
            continue
        seen.add(v)
        pptx = _TEMPLATE_DIR / f"{v}.pptx"
        mapping = _TEMPLATE_DIR / f"{v}.mapping.json"
        if pptx.exists() and mapping.exists():
            return pptx, mapping, v
    # Aucun template valide — on retourne quand même les chemins v2
    # pour produire une erreur explicite côté caller.
    return (
        _TEMPLATE_DIR / "horizon_v2.pptx",
        _TEMPLATE_DIR / "horizon_v2.mapping.json",
        "horizon_v2",
    )


# ─── Catalogue des rénovations ──────────────────────────────────────


# Catalogue de rénovations exposé au wizard frontend. Le service
# substitue dans les 16 slots `reno_item_0` ... `reno_item_15` de la
# slide 8 (FINANCES OPTIMISATION) par les items cochés. Les slots
# non utilisés sont vidés (string vide).
RENOVATIONS_CATALOGUE: List[str] = [
    "Fondation",
    "Brique",
    "Portes/fenêtres",
    "Balcons/escaliers",
    "Sous-sol",
    "Toit/sous-toit",
    "Drain entrée extérieur",
    "Chauffe-eau",
    "Panneau électrique à vérifier et corriger",
    "Finir appartement",
    "Rénovations intérieures générales",
    "Entretien extérieur général",
    "Détecteur de fumée et monoxydes",
    "Vermines",
    "Rajout d'un mur",
    "Conversion chauffage",
    "Insonorisation",
    "Plomberie majeure",
    "Électricité majeure",
    "Cuisine complète",
    "Salle de bain complète",
]


# Slots disponibles dans le template v2 (16 lignes du tableau renos)
_RENO_SLOT_COUNT = 16


# Placeholders à remplacer dans le template v2 — mêmes que dans le
# mapping JSON. Si le mapping change, ces valeurs doivent suivre.
_RENO_TEMPLATE_PLACEHOLDERS = [
    "Fondation",
    "Brique",
    "Portes et fenêtres",
    "Balcons et escaliers",
    "Sous-sol",
    "Toit et sous-toit",
    "Drain entrée extérieur",
    "Chauffe-eau",
    "Panneau électrique à vérifier et corriger",
    "Finir appartement 5 1/2 ",
    "Finir appartement 7 1/2",
    "Rénovations intérieures générales",
    "Entretien extérieur général",
    "Détecteur de fumée et monoxydes",
    "Vermines",
    "Rajout d'un mur",
]


# ─── Formatters ─────────────────────────────────────────────────────


def _money_long(n: Optional[float]) -> str:
    """Ex: `1 200 000$`."""
    if n is None:
        return "—"
    rounded = int(round(n))
    sign = "-" if rounded < 0 else ""
    abs_str = str(abs(rounded))
    with_sep = re.sub(r"\B(?=(\d{3})+(?!\d))", " ", abs_str)
    return f"{sign}{with_sep}$"


def _money_long_signed(n: Optional[float]) -> str:
    """Ex: `-65 000$` ou `347 373$` (format unifié sans espace de tête).

    Bug fix v5a : avant on produisait ` -65 000$` avec un espace devant —
    incohérent avec les autres cellules positives (`347 373$`). On aligne
    maintenant tous les montants signés sur le même format compact."""
    if n is None:
        return "—"
    rounded = int(round(n))
    if rounded < 0:
        abs_str = str(abs(rounded))
        with_sep = re.sub(r"\B(?=(\d{3})+(?!\d))", " ", abs_str)
        return f"-{with_sep}$"
    return _money_long(n)


def _money_long_with_space(n: Optional[float]) -> str:
    """Ex: `2 449 959 $` (espace avant $)."""
    return _money_long(n)[:-1] + " $"


def _money_long_lead_space(n: Optional[float]) -> str:
    """Ex: ` 2 465 059$` (espace devant — comme dans le template)."""
    return " " + _money_long(n)


def _money_short(n: Optional[float]) -> str:
    """Ex: `1800$`."""
    if n is None:
        return "—"
    return f"{int(round(n))}$"


def _money_an(n: Optional[float]) -> str:
    """Ex: `91 500 $/an`."""
    if n is None:
        return "—"
    rounded = int(round(n))
    abs_str = str(abs(rounded))
    with_sep = re.sub(r"\B(?=(\d{3})+(?!\d))", " ", abs_str)
    sign = "-" if rounded < 0 else ""
    return f"{sign}{with_sep} $/an"


def _money_plus(n: Optional[float]) -> str:
    """Ex: `+ 497 100$`."""
    if n is None:
        return "—"
    return f"+ {_money_long(n)}"


def _money_short_M(n: Optional[float]) -> str:
    """Ex: `1,2 M$` (1 décimale) pour les montants en millions."""
    if n is None:
        return "—"
    val_m = n / 1_000_000.0
    s = f"{val_m:.1f}".rstrip("0").rstrip(".")
    if not s or s == "-":
        s = "0"
    return f"{s.replace('.', ',')} M$"


def _money_approx_M(n: Optional[float]) -> str:
    """Ex: `~1,68 M$` (2 décimales)."""
    if n is None:
        return "—"
    val_m = n / 1_000_000.0
    return f"~{val_m:.2f}".replace(".", ",") + " M$"


def _money_with_delta(value: Optional[float], delta: Optional[float]) -> str:
    """Ex: `2 882 305$  (+1 685 305$)`."""
    if value is None:
        return "—"
    base = _money_long(value)
    if delta is None:
        return base
    sign = "+" if delta >= 0 else ""
    return f"{base}  ({sign}{_money_long(delta)})"


def _money_payed(value: Optional[float], paid_M: Optional[float]) -> str:
    """Ex: `1 697 100$ (payé 1.2M)`."""
    base = _money_long(value)
    if paid_M is None:
        return base
    paid_s = f"{paid_M / 1_000_000.0:.1f}".rstrip("0").rstrip(".")
    return f"{base} (payé {paid_s}M)"


def _money_payed_or_annotated(
    value: Optional[float],
    paid_M: Optional[float],
    annotation: Optional[str],
) -> str:
    """Si annotation fournie (saisie humaine), on l'utilise telle quelle
    après la valeur. Sinon, format par défaut '{val} (payé X.XM)'."""
    if annotation and annotation.strip():
        base = _money_long(value) if value else ""
        ann = annotation.strip()
        return f"{base} {ann}".strip() if base else ann
    return _money_payed(value, paid_M)


def _percent_approx(n: Optional[float]) -> str:
    """Ex: `~121 %`."""
    if n is None:
        return "—"
    return f"~{int(round(n))} %"


def _percent_2(n: Optional[float]) -> str:
    """Ex: `8,00%`."""
    if n is None:
        return "—"
    return f"{n:.2f}".replace(".", ",") + "%"


def _percent_2_fr(n: Optional[float]) -> str:
    """Ex: `3,75%`."""
    if n is None:
        return "—"
    s = f"{n:.2f}".rstrip("0").rstrip(".")
    return f"{s.replace('.', ',')}%"


def _percent_0(n: Optional[float]) -> str:
    """Ex: `100%`."""
    if n is None:
        return "—"
    return f"{int(round(n))}%"


# ─── Input dataclass ────────────────────────────────────────────────


@dataclass
class ValueAddStrategy:
    """Inputs humains du wizard (tagline + bullets + flags value-add)."""

    tagline_cover: str = ""
    qualificatif_projet: str = ""
    bullet_opp_1: str = ""
    bullet_opp_2: str = ""
    levier_principal_phrase: str = ""
    bullet_opp_3: str = ""  # Backward-compat (PR #532) — mappé sur levier_principal
    bullet_opp_4: str = ""

    # Gain potentiel callout (slide 4)
    gain_potentiel_callout: str = ""
    gain_potentiel_auto: bool = True

    # Lien Centris comparable (slide 4 — texte « Similaire en vente »)
    lien_centris_comparable: str = ""

    # Valeur du comparable Centris (slide 4 — alimente Chart 85
    # « Profit à l'achat » : catégorie « Valeur réelle » + différence
    # rouge avec le prix payé). Si vide, fallback evaluation_municipale.
    valeur_comparable_centris: float = 0.0

    # Annotation valeur marchande (slide 5)
    valeur_marchande_annotation: str = ""

    # Phase 2 du Gantt (slide 6)
    phase2_label: str = ""

    # Slide 3 — Présentation du projet (nouveau, v3) :
    # Saisie manuelle car varie (2-5) et pas extrait Centris.
    nb_etages: Optional[int] = None

    # Slide 6 — Échéancier : 5 dates éditables (M1.1, M1.2, M2.1, M2.4,
    # M3.1). Auto-suggérées selon stratégies value-add, override
    # possible. Format ISO 'YYYY-MM-DD' ou vide (= utilise auto).
    date_m1_1: str = ""
    date_m1_2: str = ""
    date_m2_1: str = ""
    date_m2_4: str = ""
    date_m3_1: str = ""

    conversion_chambres: bool = False
    nb_chambres_total: int = 0
    loyer_par_chambre: float = 0.0

    conversion_chauffage: str = "aucun"  # aucun | gaz_to_elec | elec_to_thermo
    ajout_logement_type: str = ""
    ajout_logement_loyer: float = 0.0
    optimisation_loyers_std: bool = True

    programme_schl: str = "aph_50"  # aph_50 | aph_100 | aucun

    ajout_thermopompes: int = 0
    ajout_wifi: bool = False

    # Catalogue rénovations (v2)
    renovations_selectionnees: List[str] = field(default_factory=list)
    autres_renovations: str = ""

    # Backward-compat (PR #532) : ancien input texte libre
    liste_renovations: str = ""

    # Slide 11 — ROI : callout estimation long terme (saisie humaine)
    estimation_long_terme_callout: str = ""

    # Slide 12 — Tendances : callout manuel (override). Si vide, le
    # service utilise la table Zipplex (lookup par city) ou fallback.
    tendances_callout: str = ""
    tendances_callout_manuel: str = ""
    tendances_moyenne_actuelle: float = 0.0  # Override moyenne loyer quartier

    @classmethod
    def from_dict(cls, d: Optional[Dict[str, Any]]) -> "ValueAddStrategy":
        if not d:
            return cls()
        # Backward-compat : si bullet_opp_3 fourni mais pas
        # levier_principal_phrase, on utilise bullet_opp_3.
        levier = str(
            d.get("levier_principal_phrase")
            or d.get("bullet_opp_3", "")
            or ""
        )[:200]
        renovations_sel_raw = d.get("renovations_selectionnees") or []
        if isinstance(renovations_sel_raw, str):
            renovations_sel_raw = [
                x.strip()
                for x in renovations_sel_raw.split(",")
                if x.strip()
            ]
        renovations_sel = [
            str(x)[:80]
            for x in renovations_sel_raw
            if x
        ][:32]
        nb_et_raw = d.get("nb_etages")
        try:
            nb_et = int(nb_et_raw) if nb_et_raw not in (None, "", "0") else None
        except (ValueError, TypeError):
            nb_et = None
        return cls(
            tagline_cover=str(d.get("tagline_cover", "") or "")[:200],
            qualificatif_projet=str(d.get("qualificatif_projet", "") or "")[:30],
            bullet_opp_1=str(d.get("bullet_opp_1", "") or "")[:200],
            bullet_opp_2=str(d.get("bullet_opp_2", "") or "")[:200],
            levier_principal_phrase=levier,
            bullet_opp_3=levier,  # alias
            bullet_opp_4=str(d.get("bullet_opp_4", "") or "")[:200],
            gain_potentiel_callout=str(
                d.get("gain_potentiel_callout", "") or ""
            )[:30],
            gain_potentiel_auto=bool(d.get("gain_potentiel_auto", True)),
            lien_centris_comparable=str(
                d.get("lien_centris_comparable", "") or ""
            )[:500],
            valeur_comparable_centris=float(
                d.get("valeur_comparable_centris", 0) or 0
            ),
            valeur_marchande_annotation=str(
                d.get("valeur_marchande_annotation", "") or ""
            )[:80],
            phase2_label=str(d.get("phase2_label", "") or "")[:40],
            nb_etages=nb_et,
            date_m1_1=str(d.get("date_m1_1", "") or "")[:32],
            date_m1_2=str(d.get("date_m1_2", "") or "")[:32],
            date_m2_1=str(d.get("date_m2_1", "") or "")[:32],
            date_m2_4=str(d.get("date_m2_4", "") or "")[:32],
            date_m3_1=str(d.get("date_m3_1", "") or "")[:32],
            conversion_chambres=bool(d.get("conversion_chambres", False)),
            nb_chambres_total=int(d.get("nb_chambres_total", 0) or 0),
            loyer_par_chambre=float(d.get("loyer_par_chambre", 0) or 0),
            conversion_chauffage=str(
                d.get("conversion_chauffage", "aucun") or "aucun"
            ),
            ajout_logement_type=str(d.get("ajout_logement_type", "") or ""),
            ajout_logement_loyer=float(d.get("ajout_logement_loyer", 0) or 0),
            optimisation_loyers_std=bool(
                d.get("optimisation_loyers_std", True)
            ),
            programme_schl=str(d.get("programme_schl", "aph_50") or "aph_50"),
            ajout_thermopompes=int(d.get("ajout_thermopompes", 0) or 0),
            ajout_wifi=bool(d.get("ajout_wifi", False)),
            renovations_selectionnees=renovations_sel,
            autres_renovations=str(d.get("autres_renovations", "") or "")[:500],
            liste_renovations=str(d.get("liste_renovations", "") or "")[:5000],
            estimation_long_terme_callout=str(
                d.get("estimation_long_terme_callout", "") or ""
            )[:60],
            tendances_callout=str(d.get("tendances_callout", "") or "")[:20],
            tendances_callout_manuel=str(
                d.get("tendances_callout_manuel", "") or ""
            )[:20],
            tendances_moyenne_actuelle=float(
                d.get("tendances_moyenne_actuelle", 0) or 0
            ),
        )


# ─── Données Zipplex (MVP — table constante) ───────────────────────


# Mapping quartier/ville → (moyenne loyer actuel, callout) pour la
# slide TENDANCES. MVP : table en dur jusqu'à intégration Zipplex API.
# Le service utilise `analysis.city` (ou `analysis.address` fallback)
# pour lookup approximatif (substring match).
ZIPPLEX_LOOKUP: Dict[str, Tuple[float, str]] = {
    "hochelaga": (1500.0, "+900$"),
    "hochelaga-maisonneuve": (1500.0, "+900$"),
    "mercier-hochelaga-maisonneuve": (1500.0, "+900$"),
    "plateau": (1467.0, "+450$"),
    "plateau-mont-royal": (1467.0, "+450$"),
    "mile end": (1750.0, "+650$"),
    "rosemont": (1380.0, "+520$"),
    "villeray": (1320.0, "+480$"),
    "saint-michel": (1180.0, "+380$"),
    "ahuntsic": (1420.0, "+550$"),
    "verdun": (1480.0, "+600$"),
    "ville-marie": (1850.0, "+700$"),
    "centre-sud": (1620.0, "+650$"),
    "petite-patrie": (1390.0, "+490$"),
    "côte-des-neiges": (1310.0, "+440$"),
    "saint-henri": (1530.0, "+580$"),
    "longueuil": (1180.0, "+420$"),
    "laval": (1220.0, "+460$"),
    "montréal": (1500.0, "+500$"),
}


def _zipplex_lookup(city_or_neighbourhood: Optional[str]) -> Tuple[float, str]:
    """Retourne `(moyenne_loyer, callout)` pour un quartier.

    Cherche par substring case-insensitive dans `ZIPPLEX_LOOKUP`.
    Fallback : `(1500.0, "+500$")` (Montréal moyen).
    """
    if not city_or_neighbourhood:
        return 1500.0, "+500$"
    key = city_or_neighbourhood.lower().strip()
    # Match exact
    if key in ZIPPLEX_LOOKUP:
        return ZIPPLEX_LOOKUP[key]
    # Substring match (priorité aux clés les plus longues = plus
    # spécifiques)
    matches = [(k, v) for k, v in ZIPPLEX_LOOKUP.items() if k in key]
    if matches:
        matches.sort(key=lambda kv: -len(kv[0]))
        return matches[0][1]
    # Inverse : la clé contient une partie du quartier
    matches = [(k, v) for k, v in ZIPPLEX_LOOKUP.items() if key in k]
    if matches:
        matches.sort(key=lambda kv: -len(kv[0]))
        return matches[0][1]
    return 1500.0, "+500$"


# ─── Échéancier (slide 6) — durées de phase ─────────────────────────


# Mois français pour formater les dates de la slide 6 (Gantt).
_MONTHS_FR = [
    "",
    "Janvier", "Février", "Mars", "Avril", "Mai", "Juin",
    "Juillet", "Août", "Septembre", "Octobre", "Novembre", "Décembre",
]


def _fmt_date_fr(d: date) -> str:
    """Ex: `Mai 2026`."""
    return f"{_MONTHS_FR[d.month]} {d.year}"


def _add_months(d: date, months: int) -> date:
    """Avance d'un nb de mois (approche calendrier simple)."""
    total = d.month - 1 + months
    year = d.year + total // 12
    month = total % 12 + 1
    # Jour 1 — on garde la précision à la semaine pour les overrides
    return date(year, month, 1)


def _parse_iso_date(s: str) -> Optional[date]:
    """Parse 'YYYY-MM-DD' ou retourne None."""
    if not s:
        return None
    try:
        return datetime.strptime(s.strip(), "%Y-%m-%d").date()
    except (ValueError, TypeError):
        return None


def _calc_phase_duration_months(
    strat: "ValueAddStrategy",
    nb_renos: int,
) -> Dict[str, int]:
    """Calcule la durée des phases du Gantt selon les stratégies cochées.

    Retourne : {phase1_months, phase2_months, phase3_months}
        - phase1 : Acquisition (fixe ~2 mois)
        - phase2 : Optimisation (5 si chambres, 2 si chauffage, 6 si
                   >10 rénos, 3 sinon — on prend le max si combiné)
        - phase3 : Refinancement (~12 mois après fin phase 2)
    """
    p1 = 2
    p2_candidates = [3]  # défaut minimum
    if strat.conversion_chambres:
        p2_candidates.append(5)
    if strat.conversion_chauffage and strat.conversion_chauffage != "aucun":
        p2_candidates.append(2)
    if nb_renos > 10:
        p2_candidates.append(6)
    p2 = max(p2_candidates)
    p3 = 12
    return {"phase1": p1, "phase2": p2, "phase3": p3}


def _calc_echeancier_dates(
    strat: "ValueAddStrategy",
    nb_renos: int,
    start_date: Optional[date] = None,
) -> Dict[str, date]:
    """Calcule les 5 dates clés du Gantt (M1.1, M1.2, M2.1, M2.4, M3.1).

    Logique :
        - M1.1 (lettre de financement)         : start + 1 mois
        - M1.2 (passage notaire)               : start + 2 mois (= fin P1)
        - M2.1 (fin création chambres)         : début phase 2 + min(phase2_months, 2)
        - M2.4 (fin travaux / stabilisation)   : fin phase 2 = M1.2 + phase2_months
        - M3.1 (remboursement partenaires)     : M2.4 + phase3_months
    """
    if start_date is None:
        start_date = date.today()
    dur = _calc_phase_duration_months(strat, nb_renos)
    m1_1 = _add_months(start_date, 1)
    m1_2 = _add_months(start_date, dur["phase1"])
    m2_4 = _add_months(m1_2, dur["phase2"])
    # Estimation M2.1 : un peu avant la fin de la phase 2
    m2_1 = _add_months(m1_2, max(1, dur["phase2"] // 2))
    m3_1 = _add_months(m2_4, dur["phase3"])
    return {
        "m1_1": m1_1,
        "m1_2": m1_2,
        "m2_1": m2_1,
        "m2_4": m2_4,
        "m3_1": m3_1,
    }


# ─── Build context ──────────────────────────────────────────────────


def _typology_compact(typology_json: Optional[str]) -> str:
    """`{'3.5':2,'4.5':4}` → `2x3½ | 4x4½`."""
    if not typology_json:
        return ""
    try:
        d = json.loads(typology_json)
    except Exception:
        return ""
    parts = []
    for k in ("1.5", "2.5", "3.5", "4.5", "5.5", "6.5", "7.5", "8.5"):
        v = int(d.get(k, 0) or 0)
        if v > 0:
            half = "½" if k.endswith(".5") else ""
            integer = k.split(".")[0]
            parts.append(f"{v}x{integer}{half}")
    return " | ".join(parts)


def _typology_nb_vacants(typology_json: Optional[str]) -> int:
    """Nombre estimé de logements vacants — placeholder pour MVP."""
    return 0


def _full_address(rec: LeadAnalysis) -> str:
    bits: list[str] = []
    if rec.address:
        bits.append(rec.address.strip())
    if rec.city:
        bits.append(rec.city.strip())
    return ", ".join(bits) if bits else "Adresse à confirmer"


def _best_refi_scenario(results: Optional[dict]) -> Optional[dict]:
    """Retourne le scénario de refi qui a donné le best_refi (équité max)."""
    if not results:
        return None
    scenarios = results.get("scenarios") or {}
    program = (results.get("best_refi") or {}).get("program") or ""
    mapping = {
        "SCHL standard": "refi_schl",
        "SCHL Efficacité énergétique (50 pts)": "refi_aph_50",
        "SCHL Abordabilité + Efficacité (100 pts)": "refi_aph_100",
    }
    key = mapping.get(program, "refi_aph_50")
    return scenarios.get(key) or scenarios.get("refi_aph_50") or scenarios.get(
        "refi_schl"
    )


def _suggest_phase2_label(strat: ValueAddStrategy) -> str:
    """Suggère le libellé de la phase 2 du Gantt selon les leviers
    value-add cochés.

    Priorité (du plus fort au plus faible) :
        - Création chambres       (conversion_chambres)
        - Conversion chauffage    (conversion_chauffage != aucun)
        - Travaux et rénovations  (rénos cochées sans conversion)
        - Rencontres              (défaut)
    """
    if strat.phase2_label:
        return strat.phase2_label
    if strat.conversion_chambres:
        return "Création chambres"
    if strat.conversion_chauffage and strat.conversion_chauffage != "aucun":
        return "Conversion chauffage"
    if strat.renovations_selectionnees or strat.autres_renovations:
        return "Travaux et rénovations"
    return "Rencontres"


def _calc_value_add(
    rec: LeadAnalysis, strat: ValueAddStrategy
) -> Dict[str, float]:
    """Calcule les agrégats value-add à partir des inputs humains et auto."""
    results = None
    if rec.analysis_results_json:
        try:
            results = json.loads(rec.analysis_results_json)
        except Exception:
            results = None

    nb_log = int(rec.nb_logements or 0) + (
        1 if strat.ajout_logement_type else 0
    )
    revenus_actuels = float(rec.revenus_bruts or 0)
    loyer_moyen_actuel = (
        revenus_actuels / 12.0 / nb_log
        if nb_log > 0 and revenus_actuels > 0
        else 0.0
    )

    nouveau_loyer_moyen = 0.0
    if results:
        typo = results.get("typology") or {}
        nouveau_loyer_moyen = float(typo.get("h13_loyer_pondere", 0) or 0)
        if not nouveau_loyer_moyen:
            best = _best_refi_scenario(results)
            if best:
                nouveau_loyer_moyen = float(best.get("loyer_mois", 0) or 0)

    if not nouveau_loyer_moyen:
        if strat.conversion_chambres and strat.loyer_par_chambre > 0:
            nouveau_loyer_moyen = strat.loyer_par_chambre * max(
                1, strat.nb_chambres_total // max(1, nb_log)
            )
        else:
            nouveau_loyer_moyen = loyer_moyen_actuel * 1.5

    nouveaux_revenus_an = 0.0
    if results:
        best = _best_refi_scenario(results)
        if best:
            nouveaux_revenus_an = float(best.get("revenus_totaux", 0) or 0)
    if not nouveaux_revenus_an:
        nouveaux_revenus_an = nouveau_loyer_moyen * 12.0 * nb_log

    nouvelle_valeur_marchande = 0.0
    ancienne_valeur_marchande = 0.0
    if results:
        scenarios = results.get("scenarios") or {}
        achat = scenarios.get("achat") or {}
        ancienne_valeur_marchande = float(achat.get("valeur_retenue", 0) or 0)
        best = _best_refi_scenario(results)
        if best:
            nouvelle_valeur_marchande = float(
                best.get("valeur_retenue", 0) or 0
            )

    delta_revenus_an = nouveaux_revenus_an - revenus_actuels
    delta_valeur = nouvelle_valeur_marchande - ancienne_valeur_marchande

    rci_pct = 0.0
    pvi_montant = 0.0
    if results:
        best_amount = float((results.get("best_refi") or {}).get("amount", 0) or 0)
        mdf_preteur_b = float(results.get("mdf_preteur_b", 0) or 0)
        if mdf_preteur_b > 0:
            rci_pct = (best_amount / mdf_preteur_b) * 100.0
        pvi_montant = delta_valeur

    # Gain potentiel callout = delta valeur marchande (auto) OU saisi humain
    gain_potentiel_value = delta_valeur

    return {
        "nb_log": nb_log,
        "loyer_moyen_actuel": loyer_moyen_actuel,
        "nouveau_loyer_moyen": nouveau_loyer_moyen,
        "revenus_actuels_an": revenus_actuels,
        "nouveaux_revenus_an": nouveaux_revenus_an,
        "ancienne_valeur_marchande": ancienne_valeur_marchande,
        "nouvelle_valeur_marchande": nouvelle_valeur_marchande,
        "delta_revenus_an": delta_revenus_an,
        "delta_valeur": delta_valeur,
        "rci_pct": rci_pct,
        "pvi_montant": pvi_montant,
        "gain_potentiel_value": gain_potentiel_value,
    }


def _equity_refi_pct(
    equity_neg: float,
    mdf_initial: float,
) -> str:
    """Cas équité négative : on affiche le % du capital du partenaire
    qui peut être remboursé (au lieu d'un montant négatif).

    Formule : pct = (mdf_initial + equity_neg) / mdf_initial
        - equity_neg est négatif (ex: -65 000$)
        - mdf_initial est l'investissement requis du partenaire (ex: 500 000$)
        - Résultat ≈ 87% → arrondi à 5% près = 85 ou 90%
    """
    if mdf_initial <= 0:
        return "—"
    remboursable = max(0.0, mdf_initial + equity_neg)
    pct = (remboursable / mdf_initial) * 100.0
    # Arrondi à 5% près pour un affichage sobre (90%, 85%, 95%)
    pct_rounded = round(pct / 5.0) * 5
    return f"Jusqu’à {int(pct_rounded)}%"


def _resolve_valeur_marchande(
    rec: LeadAnalysis,
    results: Optional[Dict[str, Any]],
    va: Dict[str, float],
) -> Tuple[float, str]:
    """Détermine la source canonique de la valeur marchande NOUVELLE
    pour la slide 5 (Plan création valeur).

    Hiérarchie (du plus précis au plus large) :
        1. `results.best_refi.scenario_choisi.valeur_economique`
           (= sortie moteur, déjà calculée)
        2. `results.scenarios.<best>.valeur_retenue` (idem)
        3. Calcul `nouveau_NOI / TGA` à partir des chiffres fiche
        4. Fallback : `evaluation_municipale × 1.5` (proxy faible)

    Retourne `(valeur, source_label)` — `source_label` est loggé pour
    audit et permet de vérifier rapidement la source utilisée.
    """
    # 1+2 — moteur d'analyse financière
    if va.get("nouvelle_valeur_marchande", 0) > 0:
        return float(va["nouvelle_valeur_marchande"]), "moteur_best_refi"
    # 3 — calcul direct (NOI projeté / TGA)
    if results:
        best = (results.get("best_refi") or {})
        revenus_net = float(best.get("revenus_net", 0) or 0)
        tga = float(rec.tga_pct or 4.0) / 100.0
        if revenus_net and tga > 0:
            val = revenus_net / tga
            return val, "calcul_NOI/TGA"
    # 4 — proxy faible
    eval_muni = float(rec.evaluation_municipale or 0)
    if eval_muni > 0:
        return eval_muni * 1.5, "proxy_eval_muni*1.5"
    return 0.0, "indeterminé"


def _replace_chart_data(
    slide,
    chart_name: str,
    categories: List[str],
    series: List[Tuple[str, List[float]]],
) -> bool:
    """Remplace les données d'un Chart embedded par les valeurs fournies.

    Args:
        slide: le slide python-pptx contenant le chart.
        chart_name: nom du shape (ex: "Chart 79").
        categories: liste des catégories de l'axe X (ex: ["Avant", "Après"]).
        series: liste de tuples (nom_serie, valeurs).

    Returns:
        True si la substitution a réussi, False sinon (chart introuvable
        ou erreur — non bloquant pour la génération).
    """
    try:
        from pptx.chart.data import CategoryChartData  # type: ignore
    except Exception:  # noqa: BLE001
        return False
    try:
        for shape in slide.shapes:
            if (
                hasattr(shape, "has_chart")
                and shape.has_chart
                and shape.name == chart_name
            ):
                cd = CategoryChartData()
                cd.categories = categories
                for s_name, s_vals in series:
                    cd.add_series(s_name, s_vals)
                shape.chart.replace_data(cd)
                return True
    except Exception as exc:  # noqa: BLE001
        log.warning(
            "chart replacement failed for %s: %s", chart_name, exc
        )
    return False


def _build_substitutions(
    rec: LeadAnalysis,
    strat: ValueAddStrategy,
) -> List[Tuple[int, str, str, str]]:
    """Construit la liste (slide_idx, find, replace, dupe_strategy)."""
    va = _calc_value_add(rec, strat)
    results = None
    if rec.analysis_results_json:
        try:
            results = json.loads(rec.analysis_results_json)
        except Exception:
            results = None

    achat = (results or {}).get("scenarios", {}).get("achat", {}) or {}
    aph50 = (results or {}).get("scenarios", {}).get("refi_aph_50", {}) or {}
    aph100 = (results or {}).get("scenarios", {}).get("refi_aph_100", {}) or {}

    asking = float(rec.asking_price or 0)
    revenus = float(rec.revenus_bruts or 0)
    nb_log_actuel = int(rec.nb_logements or 0)

    compact = _typology_compact(rec.typology_json)
    nb_vacants = _typology_nb_vacants(rec.typology_json)

    address_full = _full_address(rec)

    achat_dep = achat.get("depenses", {}) or {}

    investissement_requis_val = float(rec.mdf_preteur_b or va.get("delta_valeur", 0) or 0)
    if not investissement_requis_val and results:
        investissement_requis_val = float(
            results.get("mdf_preteur_b", 0) or 0
        )
    fonds_necessaires_achat_val = investissement_requis_val

    frais_demarrage_total = 0.0
    if results:
        frais_demarrage_total = float(
            results.get("frais_demarrage_total", 0) or 0
        )

    pret_max_achat = float(achat.get("financement", 0) or asking * 0.75)
    mdf_achat = max(0.0, asking - pret_max_achat) if asking else 0.0

    # APH 50/100 specifics
    revenus_nets_aph50 = float(aph50.get("revenus_net", 0) or 0)
    revenus_nets_aph100 = float(aph100.get("revenus_net", 0) or 0)
    valeur_eco_aph50 = float(aph50.get("valeur_retenue", 0) or 0)
    valeur_eco_aph100 = float(aph100.get("valeur_retenue", 0) or 0)
    pret_max_aph50 = float(aph50.get("financement", 0) or 0)
    pret_max_aph100 = float(aph100.get("financement", 0) or 0)
    # Équité dégagée : SIGNÉ pour la cellule (peut être négatif → 5271)
    equite_aph50_signed = pret_max_aph50 - pret_max_achat
    equite_aph100_signed = pret_max_aph100 - pret_max_achat
    depenses_operations_refi = float(aph50.get("revenus_totaux", 0) or 0) - revenus_nets_aph50

    # Logique conditionnelle équité refi du best
    best_scenario = _best_refi_scenario(results) or {}
    best_pret_max = float(best_scenario.get("financement", 0) or pret_max_aph50)
    best_equite = best_pret_max - pret_max_achat if asking else 0.0
    is_equity_negative = best_equite < 0
    if is_equity_negative:
        label_callout_refi = "Remboursement du partenaire"
        # TextBox 18 affiche "Jusqu'à\n{valeur}" : on remplace juste le
        # montant par un pourcentage. Le préfixe "Jusqu'à" reste dans
        # le template (paragraphe séparé).
        valeur_callout_refi_tb18 = _equity_refi_pct(
            best_equite, investissement_requis_val
        ).replace("Jusqu’à ", "")
    else:
        label_callout_refi = "Équité dégagée"
        valeur_callout_refi_tb18 = _money_long(max(0.0, best_equite))

    # WiFi / Thermopompes lines for slide 9 (refi)
    if strat.ajout_wifi or rec.ajout_wifi:
        ligne_wifi = "OUI"
    else:
        ligne_wifi = "NON"
    if strat.ajout_thermopompes or (rec.nb_thermopompes_ajoutees or 0):
        ligne_thermopompes = f"OUI ({strat.ajout_thermopompes or rec.nb_thermopompes_ajoutees})"
    else:
        ligne_thermopompes = "NON"

    reduction_energie = float(
        strat.conversion_chauffage == "elec_to_thermo" and 100
        or (rec.reduction_energie_pct or 0)
    )

    # Bourse projection (S&P 10% annualisé)
    valeur_annee_2 = investissement_requis_val * (1.10 ** 2)
    valeur_annee_10 = investissement_requis_val * (1.10 ** 10)
    equite_an2_court = pret_max_aph50 - asking if asking else 0

    # Tagline + bullets fallbacks
    quartier = (rec.city or "le secteur").strip()
    qualificatif = (strat.qualificatif_projet or "solide").strip() or "solide"
    tagline = (
        strat.tagline_cover
        or f"De l'achat au refinancement : un projet immobilier {qualificatif} à {quartier}"
    )
    bullets = [
        strat.bullet_opp_1
        or (
            f"Offre d'achat acceptée sous la valeur municipale"
            if rec.evaluation_municipale and asking < float(rec.evaluation_municipale)
            else "Acquisition à fort potentiel d'optimisation"
        ),
        strat.bullet_opp_2
        or f"Loyers moyens actuels {_money_short(va['loyer_moyen_actuel'])} | "
        f"Cible {_money_short(va['nouveau_loyer_moyen'])}",
        strat.levier_principal_phrase
        or "Marge significative de création de valeur via optimisation",
        strat.bullet_opp_4
        or f"Demande forte | Secteur {quartier}",
    ]

    # Gain potentiel callout
    if strat.gain_potentiel_auto or not strat.gain_potentiel_callout:
        gain_callout = _money_plus(va.get("gain_potentiel_value", 0))
    else:
        gain_callout = strat.gain_potentiel_callout

    # Phrase d'immeuble
    phrase_immeuble = f"Immeuble de {nb_log_actuel} logements" if nb_log_actuel else "Immeuble multilogements"
    nb_logements_phrase = f"{nb_log_actuel} logements" if nb_log_actuel else "—"
    if compact:
        phrase_typologie = f"{nb_log_actuel} logements\n {compact}\n{nb_vacants} vacants"
    else:
        phrase_typologie = f"{nb_log_actuel} logements"

    paid_M = (asking / 1_000_000.0) if asking else None

    # Phase 2 label (Gantt slide 5)
    phase2_label = _suggest_phase2_label(strat)

    # Taxes scolaires conditional label (slide 7)
    taxes_scolaires_val = float(rec.taxes_scolaires or 0)
    if taxes_scolaires_val <= 0:
        taxes_municipales_label = "Taxes municipales (inclus scolaire)"
    else:
        taxes_municipales_label = "Taxes municipales"

    # Slide 3 — auto-calculs présentation projet
    # Frais énergétiques : si énergie > 2000$/an → propriétaire paie
    energie_an = float(rec.energie or 0)
    frais_energetiques_label = (
        "Propriétaire" if energie_an > 2000.0 else "Locataire"
    )
    # Stationnements
    nb_stationnements = int(rec.nb_stationnements or 0)
    if nb_stationnements <= 0:
        stationnements_label = "Aucun"
    else:
        suffix = "places" if nb_stationnements > 1 else "place"
        stationnements_label = f"{nb_stationnements} {suffix}"
    # Année construction
    annee_construction = (
        str(rec.annee_construction)
        if rec.annee_construction
        else "À déterminer"
    )
    # Superficie habitable — pas de champ direct dans le modèle, on
    # tente `superficie_batiment` × nb_etages estimé, sinon « À
    # déterminer ». Le mapping JSON v3 expose `superficie_habitable`
    # comme variable hybrid, mais le modèle actuel n'a pas ce champ —
    # fallback explicite.
    superficie_habitable = "À déterminer"
    if rec.superficie_batiment:
        try:
            superficie_habitable = f"{int(round(float(rec.superficie_batiment)))} pi²"
        except (ValueError, TypeError):
            pass
    # Nb étages (saisie wizard manuelle — défaut 3 si inconnu)
    nb_etages = strat.nb_etages if strat.nb_etages and strat.nb_etages > 0 else 3
    nb_etages_str = str(nb_etages)

    # Slide 6 — dates échéancier auto-suggérées + overrides wizard
    nb_renos_actifs = len(strat.renovations_selectionnees or [])
    if strat.autres_renovations:
        nb_renos_actifs += len(
            [l for l in strat.autres_renovations.split("\n") if l.strip()]
        )
    auto_dates = _calc_echeancier_dates(strat, nb_renos_actifs)
    # Apply overrides (parseable ISO dates)
    final_dates: Dict[str, date] = {}
    for key, attr in [
        ("m1_1", "date_m1_1"),
        ("m1_2", "date_m1_2"),
        ("m2_1", "date_m2_1"),
        ("m2_4", "date_m2_4"),
        ("m3_1", "date_m3_1"),
    ]:
        override = _parse_iso_date(getattr(strat, attr, ""))
        final_dates[key] = override or auto_dates[key]

    # Tendances slide 11
    tendances_secteur = (rec.city or "").strip() or "le secteur"
    tendances_source_label = (
        f"Zipplex, données Q1 2026 — secteur {tendances_secteur}"
    )
    # Lookup Zipplex (ou override manuel)
    zip_moy, zip_callout = _zipplex_lookup(tendances_secteur)
    # Priorité : manuel > legacy `tendances_callout` (PR #533) > Zipplex
    tendances_callout_final = (
        strat.tendances_callout_manuel
        or strat.tendances_callout
        or zip_callout
    )
    tendances_moy_actuelle = (
        strat.tendances_moyenne_actuelle
        if strat.tendances_moyenne_actuelle > 0
        else zip_moy
    )
    # Titre dynamique slide 11
    tendances_titre = f"Tendances — {tendances_secteur}".upper()

    # Valeur marchande NOUVELLE — source canonique documentée (slide 5)
    new_val_marchande, vm_source = _resolve_valeur_marchande(
        rec, results, va
    )
    if new_val_marchande > 0:
        va["nouvelle_valeur_marchande"] = new_val_marchande
        va["delta_valeur"] = new_val_marchande - va.get(
            "ancienne_valeur_marchande", 0
        )
    log.info(
        "Valeur marchande nouvelle = %s (source=%s)",
        _money_long(new_val_marchande),
        vm_source,
    )

    # Slide 8 — chiffres frais détaillés depuis la fiche
    frais_dev = float(rec.frais_developpement or 0)
    frais_neg = float(rec.frais_negociations or 0)
    frais_dev_total = frais_dev + frais_neg
    travaux_estimes = float(rec.travaux_estimes or 0)
    # Frais autres : depuis les overrides du moteur s'ils existent
    frais_autres_overrides: Dict[str, float] = {}
    if rec.frais_demarrage_overrides_json:
        try:
            raw = json.loads(rec.frais_demarrage_overrides_json) or {}
            frais_autres_overrides = {
                k: float(v or 0) for k, v in raw.items()
            }
        except Exception:  # noqa: BLE001
            frais_autres_overrides = {}
    # Taxes bienvenue : déjà calculé par moteur (results.frais_demarrage)
    frais_demarrage_breakdown = (results or {}).get(
        "frais_demarrage_breakdown"
    ) or {}
    taxes_bienvenue = float(frais_demarrage_breakdown.get("taxes_mutation", 0) or 0)
    if not taxes_bienvenue:
        taxes_bienvenue = float(frais_autres_overrides.get("taxes_mutation", 0) or 0)
    # Intérêt revenus pendant projet
    duree_proj = int(rec.duree_projet_annees or 1)
    taux_b = float(rec.taux_interet_preteur_b_projet_pct or 8.0) / 100.0
    mdf_b = float(rec.mdf_preteur_b or 0)
    interet_revenus = mdf_b * taux_b * duree_proj if mdf_b else 0.0
    if not interet_revenus:
        interet_revenus = float(
            (results or {}).get("interet_revenus_projet", 0) or 0
        )
    # Frais évaluateur/inspection/notaire : depuis overrides ou défauts
    frais_evaluateur = float(frais_autres_overrides.get("evaluateur", 2000) or 2000)
    frais_evaluateur_2 = float(frais_autres_overrides.get("evaluateur_2", 2000) or 2000)
    frais_inspection = float(frais_autres_overrides.get("inspection", 3000) or 3000)
    frais_notaire = float(frais_autres_overrides.get("notaire", 2000) or 2000)
    frais_notaire_2 = float(frais_autres_overrides.get("notaire_2", 2000) or 2000)
    frais_avocat = float(frais_autres_overrides.get("avocat", 5000) or 5000)
    frais_courtier_1 = float(frais_autres_overrides.get("courtier_hypothecaire_1", 12000) or 12000)
    frais_courtier_2 = float(frais_autres_overrides.get("courtier_hypothecaire_2", 25000) or 25000)
    frais_rapport_eff = float(frais_autres_overrides.get("rapport_efficacite", 5000) or 5000)
    # Total frais autres (colonnes 6+7 du tableau slide 8)
    frais_autres_total = (
        frais_courtier_1 + frais_courtier_2 + taxes_bienvenue
        + frais_evaluateur + frais_evaluateur_2 + frais_inspection
        + frais_avocat + frais_notaire + frais_notaire_2
        + frais_rapport_eff + interet_revenus
    )

    # Catalogue rénovations slide 8 — on substitue les 16 placeholders
    # par les items cochés. Les slots non utilisés sont vidés.
    renos_finaux: List[str] = list(strat.renovations_selectionnees or [])
    # Backward-compat : si l'utilisateur a saisi le texte libre legacy
    # (PR #532), on essaie de splitter par newline et fusionner.
    if not renos_finaux and strat.liste_renovations:
        renos_finaux = [
            l.strip()
            for l in strat.liste_renovations.split("\n")
            if l.strip()
        ]
    # Ajouter le free-text « autres rénovations » (split par newline)
    if strat.autres_renovations:
        for line in strat.autres_renovations.split("\n"):
            line = line.strip()
            if line:
                renos_finaux.append(line)
    # Padding ou troncature aux 16 slots disponibles
    while len(renos_finaux) < _RENO_SLOT_COUNT:
        renos_finaux.append("")
    renos_finaux = renos_finaux[:_RENO_SLOT_COUNT]

    out: List[Tuple[int, str, str, str]] = [
        # Slide 0
        (0, "De l’achat au refinancement : un projet immobilier solide à Hochelaga-Maisonneuve", tagline, "first"),
        (0, "1660 Rue Saint-Clément, Montréal", address_full, "all"),
        # Slide 1
        (1, "1660 Rue Saint-Clément, Montréal", address_full, "first"),
        (1, "Immeuble de 8 logements", phrase_immeuble, "first"),
        (1, "675 000$", _money_long(investissement_requis_val), "first"),
        (1, "~121 %", _percent_approx(va["rci_pct"]), "first"),
        (1, "~1,68 M$", _money_approx_M(va["pvi_montant"]), "first"),
        # Slide 2 — Présentation du projet (Phil slide 3)
        # ⚠️ Substitutions ambiguës retirées (prix d'acquisition,
        # phrase typologie, nb_logements_phrase) — elles écrasaient
        # le format des runs (taille 14 pas gras) car la sub se faisait
        # au niveau text_frame.text fallback. Ces cellules sont
        # désormais substituées par accès direct (voir la section
        # « Substitutions directes cellules slides 3-6 » à la fin
        # de generate_offre_investissement_pptx).
        (2, "91 500 $/an", _money_an(revenus), "first"),
        (2, "171 600 $/an", _money_an(va["nouveaux_revenus_an"]), "first"),
        # Champs auto présentation projet (v3) — substitution sûre
        # via accès cellule pour superficie/annee_construction.
        # Frais énergétiques : substitution du libellé "Locataires (sauf 1)"
        (2, "Locataires (sauf 1)", frais_energetiques_label, "first"),
        # Stationnements : substitution du libellé "Aucun"
        # (premier match — vu que le cell stationnements est le dernier)
        (2, "Aucun", stationnements_label, "first"),
        # Slide 3 — Opportunité unique (Phil slide 4)
        # ⚠️ Bullet 3 (paragraphe 2 du TextBox 30) retiré ici car
        # la sub par texte écrasait le format (bold/jaune) ; à la
        # place, on remplace directement le paragraphe 2 en
        # préservant les autres bullets (voir section « Substitutions
        # directes » à la fin).
        (3, "Offre d’achat acceptée 29% sous la valeur municipale", bullets[0], "first"),
        (3, "Loyers moyens actuels 953$ | Marché 1800$", bullets[1], "first"),
        (3, "Demande pour chambres en forte hausse | Secteur Hochelaga", bullets[3], "first"),
        (3, "+ 497 100$", gain_callout, "first"),
        # Slide 3 — Lien Centris comparable (slide 4 user, idx 3)
        (
            3,
            "https://www.centris.ca/fr/multifamilial~a-vendre~montreal-mercier-hochelaga-maisonneuve/25650511",
            strat.lien_centris_comparable or "https://www.centris.ca/",
            "first",
        ),
        # Slide 4 (Plan création de valeur — Phil slide 5)
        # ⚠️ TOUS les chiffres des tableaux Avant/Après sont substitués
        # par accès direct aux cellules (Table 8 = Avant, Table 12 =
        # Après) dans la section « Substitutions directes » à la fin.
        # Cela permet de :
        #   1) Préserver le format multi-runs de cell[3][1] (valeur
        #      marchande Après : "2 882 305$  " pas gras + "(+1 685 305$)"
        #      gras jaune)
        #   2) Garantir la cohérence des sources de chiffres avec
        #      slides 3 et 4 (asking_price pour "avant", scenario
        #      choisi pour "après")
        #   3) Gérer le coût énergétique conditionnel (0$ si
        #      conversion_chauffage = elec_to_thermo)
        # Slide 5 (Échéancier Gantt — Phil slide 6)
        # Phase 2 label (Table 1 cell[4][1]) : conditionnel sur les
        # leviers value-add cochés. Si conversion_chambres=True →
        # "Création chambres" (préservé). Sinon → label suggéré
        # ("Rencontres", "Conversion chauffage", etc.).
        (5, "Création chambres", phase2_label, "first"),
        # Table 12 cell[1][1] "Fin de création des chambres" :
        # conditionnel — si conversion_chambres=False, on remplace par
        # un label neutre (ou vide). Sinon, on garde le texte.
        # (Substitution faite ici uniquement si le label suggéré est
        # DIFFÉRENT de "Création chambres" — sinon noop.)
        # Slide 6 — Conditionnel « Création de chambre » : si la
        # stratégie conversion_chambres n'est PAS cochée, on remplace
        # les textes du jalon M2.1 par un libellé neutre (rencontre
        # locataires) pour ne pas afficher « Création » à tort.
        # NOTE : la sub `phase2_label` plus haut couvre déjà Table 1
        # cell[4][1] ("Création chambres" → label suggéré). Ici on
        # complète avec Table 12 cell[1][1] ("Fin de création des
        # chambres" → libellé alternatif si non coché).
        (
            5,
            "Fin de création des chambres",
            (
                "Fin de création des chambres"
                if strat.conversion_chambres
                else "Rencontres avec locataires"
            ),
            "first",
        ),
        # Dates des jalons (M1.1, M1.2, M2.1, M2.4, M3.1) — substitue
        # les placeholders littéraux des Tables 6/12/21. Les cellules
        # sont au format "M1.1 – Juillet 2026" (espace + en-dash + espace).
        # On remplace la chaîne complète pour préserver le style.
        (5, "M1.1 – Juillet 2026", f"M1.1 – {_fmt_date_fr(final_dates['m1_1'])}", "first"),
        (5, "M1.2 – Août 2026", f"M1.2 – {_fmt_date_fr(final_dates['m1_2'])}", "first"),
        (5, "M2.1 – Oct 2026", f"M2.1 – {_fmt_date_fr(final_dates['m2_1'])}", "first"),
        (5, "M2.4 – Juillet 2028", f"M2.4 – {_fmt_date_fr(final_dates['m2_4'])}", "first"),
        (5, "M3.1 – Juillet 2028", f"M3.1 – {_fmt_date_fr(final_dates['m3_1'])}", "first"),
        # Slide 7 (FINANCES ACHAT)
        (7, "1 200 000$", _money_long(asking), "all"),
        (7, "375 000$", _money_long(frais_demarrage_total), "all"),
        (7, "91 500$", _money_long(revenus), "first"),
        (7, "953$", _money_short(va["loyer_moyen_actuel"]), "first"),
        (7, "8,00%", _percent_2(rec.taux_interet_preteur_b_projet_pct), "first"),
        (7, "2 745$", _money_long(achat_dep.get("inoccupation")), "first"),
        (7, "Taxes municipales (inclus scolaire)", taxes_municipales_label, "first"),
        (7, "12 086$", _money_long(achat_dep.get("taxes_municipales")), "first"),
        (7, "7 000$", _money_long(achat_dep.get("assurances")), "first"),
        (7, "2000$", _money_long(achat_dep.get("energie")), "first"),
        (7, "3 889$", _money_long(achat_dep.get("gestion")), "first"),
        (7, "4 880$", _money_long(achat_dep.get("entretien")), "first"),
        (7, "1 720$", _money_long(achat_dep.get("concierge")), "first"),
        (7, "58 800$", _money_long(achat.get("revenus_net")), "first"),
        (7, "900 000$", _money_long(pret_max_achat), "all"),
        (7, "300 000$", _money_long(mdf_achat), "first"),
        (7, "675 000$", _money_long(fonds_necessaires_achat_val), "all"),
        # Slide 9 (REFINANCEMENT)
        (9, "41 997$", _money_long(depenses_operations_refi), "first"),
        (9, "1788$", _money_short(va["nouveau_loyer_moyen"]), "first"),
        (9, "171 600$", _money_long(va["nouveaux_revenus_an"]), "first"),
        (9, "OUI – 1920$", ligne_wifi, "first"),
        (9, "OUI (2) – 380$", ligne_thermopompes, "first"),
        (9, "100%", _percent_0(reduction_energie), "first"),
        (9, "3,75%", _percent_2_fr(rec.taux_interet_refi_pct), "first"),
        (9, "129 603$", _money_long(revenus_nets_aph50), "first"),
        (9, "113 576$", _money_long(revenus_nets_aph100), "first"),
        (9, "2 882 305$", _money_long(valeur_eco_aph50), "all"),
        (9, " 2 465 059$", _money_long_lead_space(valeur_eco_aph100), "first"),
        (9, "2 449 959 $", _money_long_with_space(pret_max_aph50), "first"),
        (9, "2 449 959$", _money_long(pret_max_aph50), "all"),
        (9, "2 341 806$", _money_long(pret_max_aph100), "first"),
        # Équité cellules + TextBox 18 callout :
        # ⚠️ Retiré v5a : les substitutions string « 347 373$ » et
        # « 239 220$ » sont faites par accès direct cellule à la fin
        # de generate_offre_investissement_pptx (logique cohérente :
        # si équité < 0 → cellule + callout TB18 affichent un %, sinon
        # le montant). Le TextBox 18 reste géré ici via « 347 373$ »
        # car il est le SEUL textbox à contenir cette valeur après
        # substitution directe des cellules.
        (9, "347 373$", valeur_callout_refi_tb18, "first"),  # TextBox 18 callout
        # Libellé du callout (TextBox 29 + en-tête de cellule "Équité dégagée")
        (9, "Équité dégagée", label_callout_refi, "all"),
        # Slide 10 (RCI/PVI — Bourse vs Horizon)
        (10, "675 000$", _money_long(investissement_requis_val), "all"),
        (10, "817 000$", _money_long(valeur_annee_2), "first"),
        (10, "1 750 000$", _money_long(valeur_annee_10), "first"),
        (10, "780 000$", _money_long(equite_an2_court), "first"),
        # Slide 10 — Estimation long terme callout (Speech Bubble Oval 84)
        (
            10,
            "Estimations à plus de 3 M$!",
            strat.estimation_long_terme_callout
            or "Estimations à plus de 3 M$!",
            "first",
        ),
        # Slide 11 (TENDANCES) — titre dynamique + source + callout
        (11, "TENDANCES", tendances_titre, "first"),
        (
            11,
            "Zipplex, données Q1 2026 — secteur Mercier-Hochelaga-Maisonneuve",
            tendances_source_label,
            "first",
        ),
        (11, "+900$", tendances_callout_final, "first"),
    ]

    # Slide 8 — Frais détaillés (colonnes 3-4 + 6-7 du tableau renos)
    # Les chiffres du template sont substitués par ceux de la fiche.
    # ORDRE IMPORTANT : on substitue d'abord les TOTAUX (qui sont des
    # valeurs uniques) PUIS les valeurs intermédiaires — sinon une sub
    # intermédiaire peut produire une valeur identique à un total et
    # contaminer le matching suivant.
    out.extend([
        # Total rénovations : aligner avec `travaux_estimes` de la fiche
        # (correction du bug Phil — le total affichait 600 000$ du
        # template 1660 alors que la fiche dit autre chose)
        (8, "600 000$", _money_long(travaux_estimes) if travaux_estimes else "À déterminer", "all"),
        # Total frais dev (colonne 3-4)
        (8, "200 000$", _money_long(frais_dev_total), "first"),
        # Total frais autres (colonne 6-7, ligne 13)
        (8, "108 000$", _money_long(frais_autres_total), "first"),
        # Bloc « Frais de démarrage » (TextBox 8) — total grand bandeau
        (8, "375 000$", _money_long(frais_demarrage_total), "all"),
        # Détails colonne « Développement de projet »
        (8, "120 000$", _money_long(frais_dev), "first"),
        (8, "80 000$", _money_long(frais_neg), "first"),
        # Détails colonne « Frais autres »
        (8, "12 000$", _money_long(frais_courtier_1), "first"),
        (8, "25 000$", _money_long(frais_courtier_2), "first"),
        (8, "19 869$", _money_long(taxes_bienvenue), "first"),
        (8, "30 000$", _money_long(interet_revenus), "first"),
        # Évaluateur / Inspection / Avocat / Notaire / Rapport efficacité :
        # ces valeurs SONT substituées (v5a) via accès direct cellule
        # à la fin de generate_offre_investissement_pptx — la sub par
        # texte est trop ambiguë (plusieurs cellules ont 2000$ / 5000$),
        # donc on cible chaque cellule explicitement par (row, col).
    ])

    # Slide 8 : catalogue rénovations — substitue les 16 placeholders
    # par les items cochés (ou vide string si rien).
    for i, placeholder in enumerate(_RENO_TEMPLATE_PLACEHOLDERS):
        replacement = renos_finaux[i] if i < len(renos_finaux) else ""
        out.append((8, placeholder, replacement, "first"))

    return out


# ─── PPTX manipulation ──────────────────────────────────────────────


def _safe_replace_in_text_frame(text_frame, find: str, replace: str) -> int:
    """Replace `find` by `replace` inside a text frame while preserving runs."""
    if not find or find == replace:
        return 0
    count = 0
    for para in text_frame.paragraphs:
        for run in para.runs:
            if find in run.text:
                run.text = run.text.replace(find, replace)
                count += 1
                return count
    for para in text_frame.paragraphs:
        full_text = "".join(r.text for r in para.runs)
        if find in full_text:
            new_full = full_text.replace(find, replace, 1)
            if para.runs:
                para.runs[0].text = new_full
                for r in para.runs[1:]:
                    r.text = ""
                count += 1
                return count
    full = text_frame.text
    if find in full:
        new_full = full.replace(find, replace, 1)
        text_frame.text = new_full
        return 1
    return 0


def _replace_in_slide(slide, find: str, replace: str, dupe_strategy: str) -> int:
    """Walk all shapes in slide and replace."""
    if not find:
        return 0
    total = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            n = _safe_replace_in_text_frame(shape.text_frame, find, replace)
            total += n
            if total > 0 and dupe_strategy == "first":
                return total
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    n = _safe_replace_in_text_frame(cell.text_frame, find, replace)
                    total += n
                    if total > 0 and dupe_strategy == "first":
                        return total
    return total


def _replace_table_cell_text(
    slide,
    table_name: str,
    row_idx: int,
    col_idx: int,
    new_text: str,
) -> bool:
    """Remplace le texte d'une cellule de table par accès direct.

    Préserve le formatting du premier run + supprime les autres runs/
    paragraphes. Utilisé pour les substitutions où la valeur cible est
    ambiguë (ex: "3" → nb_etages dans slide 2).
    """
    try:
        for shape in slide.shapes:
            if not getattr(shape, "has_table", False):
                continue
            if shape.name != table_name:
                continue
            tbl = shape.table
            if row_idx >= len(tbl.rows):
                return False
            row = tbl.rows[row_idx]
            cells = list(row.cells)
            if col_idx >= len(cells):
                return False
            cell = cells[col_idx]
            tf = cell.text_frame
            # Préserve le 1er paragraphe + 1er run, vide le reste
            if tf.paragraphs:
                first_para = tf.paragraphs[0]
                if first_para.runs:
                    first_para.runs[0].text = new_text
                    for r in first_para.runs[1:]:
                        r.text = ""
                else:
                    first_para.text = new_text
                # Vide les paragraphes supplémentaires
                for p in tf.paragraphs[1:]:
                    for r in p.runs:
                        r.text = ""
            else:
                tf.text = new_text
            return True
    except Exception as exc:  # noqa: BLE001
        log.warning(
            "table cell replacement failed for %s[%s][%s]: %s",
            table_name, row_idx, col_idx, exc,
        )
    return False


def _replace_cell_paragraphs(
    slide,
    table_name: str,
    row_idx: int,
    col_idx: int,
    paragraphs_text: List[str],
) -> bool:
    """Remplace le contenu d'une cellule avec une LISTE de paragraphes.

    Préserve le format du premier run de CHAQUE paragraphe existant
    (police, couleur, taille, gras) et vide les runs supplémentaires.
    Si la liste contient PLUS de paragraphes que la cellule en a, on
    duplique le format du dernier paragraphe existant.
    Si elle en contient MOINS, on vide les paragraphes restants.

    Utile pour Table 12 cell[0][1] slide 3 ("8 logements\\n 2x3½...\\n2
    vacants") où l'on veut respecter le format taille 14 pas-gras des
    3 paragraphes du template sans recourir à un remplacement texte qui
    écrase tout le formatting.
    """
    try:
        for shape in slide.shapes:
            if not getattr(shape, "has_table", False):
                continue
            if shape.name != table_name:
                continue
            tbl = shape.table
            if row_idx >= len(tbl.rows):
                return False
            row = tbl.rows[row_idx]
            cells = list(row.cells)
            if col_idx >= len(cells):
                return False
            cell = cells[col_idx]
            tf = cell.text_frame
            paras = list(tf.paragraphs)
            n_target = len(paragraphs_text)
            n_existing = len(paras)
            for i in range(max(n_target, n_existing)):
                if i < n_existing:
                    para = paras[i]
                    if i < n_target:
                        # Remplace le texte dans le premier run, vide
                        # les runs suivants
                        if para.runs:
                            para.runs[0].text = paragraphs_text[i]
                            for r in para.runs[1:]:
                                r.text = ""
                        else:
                            para.text = paragraphs_text[i]
                    else:
                        # Vide le paragraphe surplus
                        for r in para.runs:
                            r.text = ""
                else:
                    # Pas assez de paragraphes existants — on ajoute
                    new_para = tf.add_paragraph()
                    new_para.text = paragraphs_text[i]
            return True
    except Exception as exc:  # noqa: BLE001
        log.warning(
            "cell paragraphs replacement failed for %s[%s][%s]: %s",
            table_name, row_idx, col_idx, exc,
        )
    return False


def _replace_textbox_paragraph(
    slide,
    shape_name: str,
    para_idx: int,
    new_text: str,
    *,
    force_bold: Optional[bool] = None,
    force_color_rgb: Optional[str] = None,
    force_size_pt: Optional[float] = None,
) -> bool:
    """Remplace un paragraphe précis d'un TextBox en préservant le format.

    Le 1er run du paragraphe est utilisé pour le nouveau texte ; les
    runs suivants sont vidés. Si `force_*` est fourni, ces attributs
    écrasent le format du run.

    Utilisé pour le bullet 3 (paragraphe idx=2) du TextBox 30 slide 4
    où l'on doit forcer bold=False, color=blanc, size=14 pour matcher
    les bullets 1/2/4.
    """
    try:
        from pptx.dml.color import RGBColor  # type: ignore
        from pptx.util import Pt  # type: ignore
    except Exception:  # noqa: BLE001
        return False
    try:
        for shape in slide.shapes:
            if shape.name != shape_name:
                continue
            if not shape.has_text_frame:
                return False
            tf = shape.text_frame
            paras = list(tf.paragraphs)
            if para_idx >= len(paras):
                return False
            para = paras[para_idx]
            if not para.runs:
                para.text = new_text
                run = para.runs[0] if para.runs else None
            else:
                para.runs[0].text = new_text
                for r in para.runs[1:]:
                    r.text = ""
                run = para.runs[0]
            if run is None:
                return True
            if force_bold is not None:
                run.font.bold = force_bold
            if force_color_rgb:
                try:
                    run.font.color.rgb = RGBColor.from_string(force_color_rgb)
                except Exception as exc:  # noqa: BLE001
                    log.warning(
                        "color set failed for %s p[%s]: %s",
                        shape_name, para_idx, exc,
                    )
            if force_size_pt is not None:
                run.font.size = Pt(force_size_pt)
            return True
    except Exception as exc:  # noqa: BLE001
        log.warning(
            "textbox paragraph replacement failed for %s p[%s]: %s",
            shape_name, para_idx, exc,
        )
    return False


def _replace_cell_keep_first_run(
    slide,
    table_name: str,
    row_idx: int,
    col_idx: int,
    new_text: str,
) -> bool:
    """Variante de `_replace_table_cell_text` qui ne touche pas au
    formatting du 1er run (gras/couleur/taille). Utile pour les
    cellules dont on veut juste changer le texte sans toucher au style.

    Différence avec `_replace_table_cell_text` : aucune modification
    du formatting du run, juste `.text = new_text`.
    """
    return _replace_table_cell_text(
        slide, table_name, row_idx, col_idx, new_text
    )


def _replace_cell_two_runs(
    slide,
    table_name: str,
    row_idx: int,
    col_idx: int,
    text_run_0: str,
    text_run_1: str,
) -> bool:
    """Remplace les 2 premiers runs du 1er paragraphe d'une cellule.

    Utilisé pour Table 12 cell[3][1] slide 5 (valeur marchande Après
    refi) : run0 = "2 882 305$  " (size 12 pas gras) + run1 =
    "(+1 685 305$)" (size 10 gras jaune). Les formats des 2 runs sont
    préservés ; on ne change que le texte.
    """
    try:
        for shape in slide.shapes:
            if not getattr(shape, "has_table", False):
                continue
            if shape.name != table_name:
                continue
            tbl = shape.table
            if row_idx >= len(tbl.rows):
                return False
            row = tbl.rows[row_idx]
            cells = list(row.cells)
            if col_idx >= len(cells):
                return False
            cell = cells[col_idx]
            tf = cell.text_frame
            if not tf.paragraphs:
                return False
            first_para = tf.paragraphs[0]
            runs = list(first_para.runs)
            if len(runs) >= 2:
                runs[0].text = text_run_0
                runs[1].text = text_run_1
                # Vider les runs suivants
                for r in runs[2:]:
                    r.text = ""
            elif len(runs) == 1:
                # Une seule run — on met les 2 textes concaténés dans le run 0
                runs[0].text = text_run_0 + text_run_1
            else:
                first_para.text = text_run_0 + text_run_1
            # Vider les paragraphes suivants
            for p in tf.paragraphs[1:]:
                for r in p.runs:
                    r.text = ""
            return True
    except Exception as exc:  # noqa: BLE001
        log.warning(
            "two-runs cell replacement failed for %s[%s][%s]: %s",
            table_name, row_idx, col_idx, exc,
        )
    return False


def _condense_empty_table_rows(
    slide,
    table_name: str,
    body_start: int,
    body_end_exclusive: int,
    col_to_check: int = 0,
) -> int:
    """Supprime les lignes vides du body d'un tableau (entre header et total).

    Utilisé pour le tableau rénos slide 9 (idx 8) : après substitution des
    16 placeholders par les items cochés, on retire les lignes dont la
    cellule libellé est vide pour que le Total reste collé aux items.

    Args:
        slide: slide python-pptx contenant la table.
        table_name: nom du shape (ex: "Table 2").
        body_start: index 0-based de la 1ère ligne body (inclusive).
        body_end_exclusive: index de la 1ère ligne après le body
            (= ligne du Total).
        col_to_check: colonne à inspecter pour considérer la ligne vide
            (par défaut 0 — libellé rénos).

    Returns:
        Nombre de lignes retirées (0 si rien — ex: shape absente).
    """
    try:
        from pptx.oxml.ns import qn  # type: ignore
    except Exception:  # noqa: BLE001
        return 0
    try:
        for shape in slide.shapes:
            if not getattr(shape, "has_table", False):
                continue
            if shape.name != table_name:
                continue
            tbl = shape.table
            tbl_xml = tbl._tbl
            rows_xml = tbl_xml.findall(qn("a:tr"))
            n_rows = len(rows_xml)
            if body_start >= n_rows or body_end_exclusive > n_rows:
                return 0
            removed = 0
            # Parcourir le body et identifier les lignes vides
            for i in range(body_start, body_end_exclusive):
                row_xml = rows_xml[i]
                # Extraire le texte de la cellule col_to_check
                # tc = a:tc enfants de a:tr
                tcs = row_xml.findall(qn("a:tc"))
                if col_to_check >= len(tcs):
                    continue
                tc = tcs[col_to_check]
                # Concat tous les a:t
                cell_text = "".join(
                    (t.text or "") for t in tc.iter(qn("a:t"))
                )
                if not cell_text.strip():
                    tbl_xml.remove(row_xml)
                    removed += 1
            return removed
    except Exception as exc:  # noqa: BLE001
        log.warning(
            "table condensation failed for %s: %s", table_name, exc
        )
    return 0


def _replace_photo(slide, shape_name: str, new_blob: bytes) -> bool:
    """Replace a Picture shape's underlying image while keeping geometry."""
    try:
        for shape in slide.shapes:
            if shape.name == shape_name and shape.shape_type == 13:  # PICTURE
                left, top, width, height = (
                    shape.left,
                    shape.top,
                    shape.width,
                    shape.height,
                )
                sp = shape._element
                sp.getparent().remove(sp)
                slide.shapes.add_picture(
                    io.BytesIO(new_blob),
                    left,
                    top,
                    width=width,
                    height=height,
                )
                return True
    except Exception as exc:
        log.warning("photo replacement failed for %s: %s", shape_name, exc)
    return False


# ─── Public API ─────────────────────────────────────────────────────


def _slugify_address(addr: Optional[str], fallback_id: int) -> str:
    if not addr or not addr.strip():
        return str(fallback_id)
    decomposed = unicodedata.normalize("NFKD", addr.strip())
    ascii_only = "".join(
        ch for ch in decomposed if not unicodedata.combining(ch)
    )
    underscored = re.sub(r"[\s\-,'’]+", "_", ascii_only)
    cleaned = re.sub(r"[^A-Za-z0-9_]", "", underscored)
    cleaned = re.sub(r"_+", "_", cleaned).strip("_")
    return cleaned or str(fallback_id)


def offre_investissement_pptx_filename(rec: LeadAnalysis) -> str:
    slug = _slugify_address(rec.address, rec.id)
    today = datetime.now().strftime("%Y-%m-%d")
    return f"Offre_Investissement_{slug}_{today}.pptx"


def get_renovations_catalogue() -> List[str]:
    """Expose le catalogue de rénovations au frontend (pour le wizard).
    Utilisé par le endpoint GET /lead-analyses/offre-investissement/catalogue."""
    return list(RENOVATIONS_CATALOGUE)


async def generate_offre_investissement_pptx(
    db: AsyncSession,
    analysis_id: int,
    value_add_strategy: Optional[Dict[str, Any]] = None,
    photos: Optional[List[bytes]] = None,
    photo_attachment_ids: Optional[List[int]] = None,
) -> Tuple[bytes, str]:
    """Generate a .pptx investment offer for the given lead analysis.

    Args:
        db: AsyncSession.
        analysis_id: LeadAnalysis id.
        value_add_strategy: human inputs from the wizard (or None for defaults).
        photos: list of raw image bytes (cover, exterieur, carte, tendances).
        photo_attachment_ids: list of LeadAnalysisAttachment ids.

    Returns:
        (pptx_bytes, template_version) : le `.pptx` généré + la version
        du template utilisée (ex: "horizon_v2", "horizon_v1").

    Raises:
        ValueError: if the analysis is not found or the template is missing.
    """
    rec = await db.get(LeadAnalysis, analysis_id)
    if rec is None:
        raise ValueError(f"LeadAnalysis {analysis_id} introuvable")

    template_path, mapping_path, template_version = _resolve_template_paths()
    if not template_path.exists():
        raise ValueError(
            f"Template introuvable : {template_path}. "
            f"Vérifiez le déploiement des assets."
        )

    from pptx import Presentation  # type: ignore

    strat = ValueAddStrategy.from_dict(value_add_strategy)

    # Resolve photos
    resolved_photos: List[bytes] = []
    if photos:
        resolved_photos.extend(photos)
    if photo_attachment_ids:
        from sqlalchemy import select

        q = select(LeadAnalysisAttachment).where(
            LeadAnalysisAttachment.id.in_(photo_attachment_ids),
            LeadAnalysisAttachment.lead_analysis_id == rec.id,
        )
        result = await db.execute(q)
        for att in result.scalars().all():
            resolved_photos.append(att.blob)

    # Open template
    prs = Presentation(str(template_path))

    # Apply text substitutions
    substitutions = _build_substitutions(rec, strat)
    applied = 0
    skipped = 0
    for slide_idx, find, replace, dupe_strategy in substitutions:
        if not find or find == replace:
            skipped += 1
            continue
        if slide_idx >= len(prs.slides):
            continue
        slide = prs.slides[slide_idx]
        n = _replace_in_slide(slide, find, replace, dupe_strategy)
        if n:
            applied += n
        else:
            skipped += 1
    log.info(
        "Offre PPTX generated (%s): %s subs applied, %s skipped (analysis %s)",
        template_version,
        applied,
        skipped,
        analysis_id,
    )

    # ─── Substitutions directes cellules slides 3-6 (v4) ──────────
    # Ces substitutions remplacent le texte de cellules précises tout
    # en PRÉSERVANT le formatage (police, taille, couleur, gras).
    # Elles complètent la phase de substitution texte ci-dessus :
    # certaines cellules ont un formatting multi-runs ou multi-
    # paragraphes que la sub texte écrase. On les traite ici.
    rec_asking = float(rec.asking_price or 0)
    rec_revenus = float(rec.revenus_bruts or 0)
    rec_eval_muni = float(rec.evaluation_municipale or 0)
    rec_energie = float(rec.energie or 0)
    rec_nb_log = int(rec.nb_logements or 0)
    _results_local = None
    if rec.analysis_results_json:
        try:
            _results_local = json.loads(rec.analysis_results_json)
        except Exception:  # noqa: BLE001
            _results_local = None
    va_local = _calc_value_add(rec, strat)
    # Force re-resolve nouvelle valeur marchande pour charts/slide 5
    new_vm_local, _vm_src = _resolve_valeur_marchande(
        rec, _results_local, va_local
    )
    if new_vm_local > 0:
        va_local["nouvelle_valeur_marchande"] = new_vm_local
        va_local["delta_valeur"] = new_vm_local - va_local.get(
            "ancienne_valeur_marchande", 0
        )
    # Coût énergétique avant/après slide 5 :
    #   - Avant = energie de la fiche (rec.energie)
    #   - Après = 0 si conversion_chauffage = elec_to_thermo, sinon energie
    cost_energie_after = (
        0.0
        if strat.conversion_chauffage == "elec_to_thermo"
        else rec_energie
    )

    # ─── Slide 3 (idx 2) — Présentation du projet ─────────────────
    if len(prs.slides) > 2:
        s2 = prs.slides[2]
        # Table 11 cell[0][1] — Prix d'acquisition (= asking_price)
        # Source garantie : rec.asking_price (= "prix demandé" fiche)
        if rec_asking > 0:
            _replace_table_cell_text(
                s2, "Table 11", 0, 1, _money_short_M(rec_asking)
            )
        # Table 11 cell[1][1] — Revenus actuels (= revenus_bruts)
        if rec_revenus > 0:
            _replace_table_cell_text(
                s2, "Table 11", 1, 1, _money_an(rec_revenus)
            )
        # Table 11 cell[2][1] — Revenus potentiels après rénos
        nouv_rev = float(va_local.get("nouveaux_revenus_an", 0) or 0)
        if nouv_rev > 0:
            _replace_table_cell_text(
                s2, "Table 11", 2, 1, _money_an(nouv_rev)
            )
        # Table 12 cell[0][1] — Nombre d'unités (multi-paragraphes,
        # format taille 14 pas gras à préserver) :
        #   p[0] = "{n} logements"
        #   p[1] = " {typology compact}"  (ex: " 2 × 3½ | 4 × 4½ ...")
        #   p[2] = "{vacants} vacants"  (si applicable)
        compact_local = _typology_compact(rec.typology_json)
        nb_vacants_local = _typology_nb_vacants(rec.typology_json)
        nb_log_phrase_local = (
            f"{rec_nb_log} logements" if rec_nb_log else "—"
        )
        paragraphs_nb_unites: List[str] = [nb_log_phrase_local]
        if compact_local:
            paragraphs_nb_unites.append(f" {compact_local}")
        if nb_vacants_local > 0:
            paragraphs_nb_unites.append(f"{nb_vacants_local} vacants")
        _replace_cell_paragraphs(
            s2, "Table 12", 0, 1, paragraphs_nb_unites
        )
        # Table 12 cell[1][1] — Nb d'étages (saisie manuelle)
        nb_et_local = (
            strat.nb_etages
            if strat.nb_etages and strat.nb_etages > 0
            else None
        )
        if nb_et_local is not None:
            _replace_table_cell_text(
                s2, "Table 12", 1, 1, str(nb_et_local)
            )

    # ─── Slide 4 (idx 3) — Opportunité unique ─────────────────────
    if len(prs.slides) > 3:
        s3 = prs.slides[3]
        # TextBox 30 paragraphe 2 = bullet 3 (« levier principal »).
        # Format forcé : bold=False, color=blanc (FFFFFF), size=14.
        # Pas de label "Levier principal:" en préfixe : c'était un
        # label du wizard, jamais inclus dans la valeur envoyée.
        levier_text = (
            strat.levier_principal_phrase
            or strat.bullet_opp_3
            or "Marge significative de création de valeur via optimisation"
        )
        _replace_textbox_paragraph(
            s3,
            "TextBox 30",
            2,
            levier_text,
            force_bold=False,
            force_color_rgb="FFFFFF",
            force_size_pt=14,
        )

        # ─── Chart 79 (LINE — VALEUR MARCHANDE) ───────────────────
        # Bug fixé : Avant = prix payé (asking_price), Après refi =
        # valeur économique du scénario choisi (programme_schl).
        scenarios_local = (_results_local or {}).get("scenarios", {}) or {}
        if strat.programme_schl == "aph_100":
            chosen_scenario = scenarios_local.get("refi_aph_100", {}) or {}
        elif strat.programme_schl == "aph_50":
            chosen_scenario = scenarios_local.get("refi_aph_50", {}) or {}
        else:
            chosen_scenario = _best_refi_scenario(_results_local) or {}
        valeur_eco_refi = float(
            chosen_scenario.get("valeur_retenue", 0)
            or chosen_scenario.get("valeur_economique", 0)
            or 0
        )
        if not valeur_eco_refi:
            valeur_eco_refi = new_vm_local
        if rec_asking > 0 and valeur_eco_refi > 0:
            _replace_chart_data(
                s3,
                "Chart 79",
                ["Avant", "Après refi"],
                [("PDM", [float(rec_asking), float(valeur_eco_refi)])],
            )

        # ─── Chart 85 (COLUMN_STACKED — PROFIT À L'ACHAT) ─────────
        # Bug fixé :
        #   - Catégories : ["Valeur d'achat", "Valeur réelle"]
        #     (au lieu de "Valeur municipale")
        #   - Valeur "Valeur réelle" = valeur_comparable_centris
        #     (nouveau champ wizard) ou fallback evaluation_municipale
        #   - 2 séries empilées :
        #       Series 0 "Prix payé" (jaune) : [prix, prix]
        #       Series 1 "Écart comparable" (rouge) : [0, delta]
        comparable_val = (
            strat.valeur_comparable_centris
            if strat.valeur_comparable_centris > 0
            else rec_eval_muni
        )
        if rec_asking > 0 and comparable_val > 0:
            delta_comparable = max(0.0, comparable_val - rec_asking)
            _replace_chart_data(
                s3,
                "Chart 85",
                ["Valeur d'achat", "Valeur réelle"],
                [
                    ("Prix payé", [rec_asking, rec_asking]),
                    ("Écart comparable", [0.0, delta_comparable]),
                ],
            )

        # TextBox 30 paragraphe 4 = "Similaire en vente à 2M$ <lien>"
        # On remplace "à 2M$" par la valeur du comparable Centris si
        # fournie (en M$ court). Évite que le deck affiche un montant
        # hardcodé du template 1660.
        if comparable_val > 0:
            # Sub texte ciblé : "à 2M$" → "à {comparable_val court}"
            comparable_str = _money_short_M(comparable_val)
            _replace_in_slide(
                s3, "à 2M$", f"à {comparable_str}", "first"
            )

    # ─── Slide 5 (idx 4) — Plan création de valeur ────────────────
    # Sources documentées (cohérence avec slides 3 et 4) :
    #   Table 8 (Avant)   : Moyenne locative, Revenus locatifs, Valeur
    #                       marchande, Coût énergétique = depuis la fiche
    #   Table 12 (Après)  : Loyer pondéré projeté, Revenus refi, Valeur
    #                       économique refi, Coût énergétique projeté
    if len(prs.slides) > 4:
        s4 = prs.slides[4]
        # ── Table 8 (AVANT) ──
        # cell[1][1] — Moyenne locative actuelle = revenus_bruts / 12 / nb_log
        loyer_av = float(va_local.get("loyer_moyen_actuel", 0) or 0)
        if loyer_av > 0:
            _replace_table_cell_text(
                s4, "Table 8", 1, 1, _money_short(loyer_av)
            )
        # cell[2][1] — Revenus locatifs annuels (avant) = revenus_bruts
        if rec_revenus > 0:
            _replace_table_cell_text(
                s4, "Table 8", 2, 1, _money_long(rec_revenus)
            )
        # cell[3][1] — Valeur marchande AVANT :
        #   Si annotation humaine (valeur_marchande_annotation) fournie,
        #     on l'utilise telle quelle (ex: "Payé 1 100 000$").
        #   Sinon, format "{valeur_calculee} (payé X.XM)".
        ancienne_vm_local = float(
            va_local.get("ancienne_valeur_marchande", 0) or 0
        )
        if strat.valeur_marchande_annotation:
            # Annotation custom — affichage tel quel
            # (ex: "Payé 1 100 000$", "(payé 1.2M)")
            vm_avant_text = strat.valeur_marchande_annotation.strip()
        elif ancienne_vm_local > 0 and rec_asking > 0:
            # Format défaut : "{valeur} (payé {prix M$})"
            paid_M_str = f"{rec_asking / 1_000_000.0:.1f}".rstrip("0").rstrip(".")
            if not paid_M_str:
                paid_M_str = "0"
            vm_avant_text = f"{_money_long(ancienne_vm_local)} (payé {paid_M_str}M)"
        elif rec_asking > 0:
            # Sans valeur municipale fiable, on affiche juste le prix payé
            vm_avant_text = f"Payé {_money_long(rec_asking)}"
        else:
            vm_avant_text = "—"
        _replace_table_cell_text(
            s4, "Table 8", 3, 1, vm_avant_text
        )
        # cell[4][1] — Coût énergétique avant = rec.energie (annuel)
        _replace_table_cell_text(
            s4, "Table 8", 4, 1, _money_long(rec_energie)
            if rec_energie > 0 else "0$"
        )
        # ── Table 12 (APRÈS) ──
        # cell[1][1] — Moyenne locative projetée
        loyer_ap = float(va_local.get("nouveau_loyer_moyen", 0) or 0)
        if loyer_ap > 0:
            _replace_table_cell_text(
                s4, "Table 12", 1, 1, _money_short(loyer_ap)
            )
        # cell[2][1] — Revenus locatifs annuels projetés
        nouv_rev_ap = float(va_local.get("nouveaux_revenus_an", 0) or 0)
        if nouv_rev_ap > 0:
            _replace_table_cell_text(
                s4, "Table 12", 2, 1, _money_long(nouv_rev_ap)
            )
        # cell[3][1] — Valeur marchande APRÈS (2 runs : montant pas
        # gras + (+delta) gras jaune)
        nouv_vm_ap = float(
            va_local.get("nouvelle_valeur_marchande", 0) or 0
        )
        delta_vm = float(va_local.get("delta_valeur", 0) or 0)
        if nouv_vm_ap > 0:
            _replace_cell_two_runs(
                s4, "Table 12", 3, 1,
                f"{_money_long(nouv_vm_ap)}  ",
                f"(+{_money_long(delta_vm)})",
            )
        # cell[4][1] — Coût énergétique projeté
        _replace_table_cell_text(
            s4, "Table 12", 4, 1,
            _money_long(cost_energie_after)
            if cost_energie_after > 0 else "0$"
        )

    # ─── Substitutions directes cellules slides 8-10 (v5a) ────────
    # PR v5a : corrige slides 8, 9, 10 (visibles = idx 7, 8, 9) avec
    # accès direct par cellule pour éviter les substitutions string
    # ambiguës (plusieurs cellules ont 2000$, 5000$, "8", etc.).
    # Les variables nécessaires (frais autres, équité) sont recalculées
    # ici à partir de rec/_results_local pour cohérence avec
    # `_build_substitutions` (DRY évitable : section additive v5a).
    _achat_local = (_results_local or {}).get("scenarios", {}).get("achat", {}) or {}
    _aph50_local = (_results_local or {}).get("scenarios", {}).get("refi_aph_50", {}) or {}
    _aph100_local = (_results_local or {}).get("scenarios", {}).get("refi_aph_100", {}) or {}
    _frais_breakdown = (_results_local or {}).get("frais_demarrage_breakdown") or {}
    _frais_overrides: Dict[str, float] = {}
    if rec.frais_demarrage_overrides_json:
        try:
            _frais_overrides = {
                k: float(v or 0)
                for k, v in (json.loads(rec.frais_demarrage_overrides_json) or {}).items()
            }
        except Exception:  # noqa: BLE001
            _frais_overrides = {}
    _taxes_bienvenue_local = float(_frais_breakdown.get("taxes_mutation", 0) or 0)
    if not _taxes_bienvenue_local:
        _taxes_bienvenue_local = float(_frais_overrides.get("taxes_mutation", 0) or 0)
    _duree_proj_local = int(rec.duree_projet_annees or 1)
    _taux_b_local = float(rec.taux_interet_preteur_b_projet_pct or 8.0) / 100.0
    _mdf_b_local = float(rec.mdf_preteur_b or 0)
    _interet_revenus_local = _mdf_b_local * _taux_b_local * _duree_proj_local if _mdf_b_local else 0.0
    if not _interet_revenus_local:
        _interet_revenus_local = float(
            (_results_local or {}).get("interet_revenus_projet", 0) or 0
        )
    # Frais autres détaillés : depuis overrides fiche ou défauts industrie
    frais_evaluateur = float(_frais_overrides.get("evaluateur", 2000) or 2000)
    frais_evaluateur_2 = float(_frais_overrides.get("evaluateur_2", 2000) or 2000)
    frais_inspection = float(_frais_overrides.get("inspection", 3000) or 3000)
    frais_avocat = float(_frais_overrides.get("avocat", 5000) or 5000)
    frais_notaire = float(_frais_overrides.get("notaire", 2000) or 2000)
    frais_notaire_2 = float(_frais_overrides.get("notaire_2", 2000) or 2000)
    frais_rapport_eff = float(_frais_overrides.get("rapport_efficacite", 5000) or 5000)
    frais_courtier_1 = float(_frais_overrides.get("courtier_hypothecaire_1", 12000) or 12000)
    frais_courtier_2 = float(_frais_overrides.get("courtier_hypothecaire_2", 25000) or 25000)
    # Total frais autres (re-calculé avec les valeurs effectivement utilisées)
    frais_autres_total = (
        frais_courtier_1 + frais_courtier_2 + _taxes_bienvenue_local
        + frais_evaluateur + frais_evaluateur_2 + frais_inspection
        + frais_avocat + frais_notaire + frais_notaire_2
        + frais_rapport_eff + _interet_revenus_local
    )
    # Équité scénarios refi pour slide 10
    _pret_max_achat_local = float(
        _achat_local.get("financement", 0) or rec_asking * 0.75
    )
    _pret_max_aph50_local = float(_aph50_local.get("financement", 0) or 0)
    _pret_max_aph100_local = float(_aph100_local.get("financement", 0) or 0)
    equite_aph50_signed = _pret_max_aph50_local - _pret_max_achat_local
    equite_aph100_signed = _pret_max_aph100_local - _pret_max_achat_local
    investissement_requis_val = float(
        rec.mdf_preteur_b
        or (_results_local or {}).get("mdf_preteur_b", 0)
        or 0
    )

    # ─── Slide 8 visible (idx 7) — Finances Achat ────────────────
    # CELL[4][1] = Nombre de logements. Le template a "8" hardcodé.
    if len(prs.slides) > 7 and rec_nb_log > 0:
        _replace_table_cell_text(
            prs.slides[7], "Table 2", 4, 1, str(rec_nb_log)
        )

    # ─── Slide 9 visible (idx 8) — Optimisation ──────────────────
    # Colonne "Frais autres" (col 7) : 7 cellules génériques (rows
    # 5-11) qui avaient le même placeholder text (2000$/3000$/5000$)
    # — impossible à différencier par sub texte. On les remplace ici
    # par cellule. Total en CELL[13][7].
    if len(prs.slides) > 8:
        s8 = prs.slides[8]
        _replace_table_cell_text(
            s8, "Table 2", 5, 7, _money_long(frais_evaluateur)
        )
        _replace_table_cell_text(
            s8, "Table 2", 6, 7, _money_long(frais_evaluateur_2)
        )
        _replace_table_cell_text(
            s8, "Table 2", 7, 7, _money_long(frais_inspection)
        )
        _replace_table_cell_text(
            s8, "Table 2", 8, 7, _money_long(frais_avocat)
        )
        _replace_table_cell_text(
            s8, "Table 2", 9, 7, _money_long(frais_notaire)
        )
        _replace_table_cell_text(
            s8, "Table 2", 10, 7, _money_long(frais_notaire_2)
        )
        _replace_table_cell_text(
            s8, "Table 2", 11, 7, _money_long(frais_rapport_eff)
        )
        # Total frais autres CELL[13][7] : re-sub directe (le total
        # est recalculé proprement avec les valeurs effectives utilisées,
        # alors que la sub texte « 108 000$ » du bloc principal risque
        # de rater si le total réel diffère du template)
        _replace_table_cell_text(
            s8, "Table 2", 13, 7, _money_long(frais_autres_total)
        )

        # Condensation lignes vides du tableau rénos.
        # Layout template : 21 rows. body rénos = rows 2..17 (16 slots),
        # row 18 = Total, row 19 = "***Certains frais financés".
        # On condense uniquement le body (rows 2..17). Les cellules
        # dont le libellé col 0 a été substitué par "" sont retirées
        # physiquement (le Total et la mention restent en bas).
        # ⚠️ DOIT être après les substitutions de placeholders (faites
        # via _replace_in_slide dans la boucle substitutions).
        _condense_empty_table_rows(
            s8, "Table 2", body_start=2, body_end_exclusive=18,
            col_to_check=0,
        )

    # ─── Slide 10 visible (idx 9) — Refinancement ────────────────
    # CELL[5][1] = Nombre de logements ("8" hardcodé dans template).
    # CELL[9][4] = Équité scénario 1 (AHP50), CELL[9][7] = scénario 2
    # (AHP100). Logique cohérente cellule + callout TextBox 18 :
    #   - équité du scénario >= 0 → afficher montant en $ (format
    #     unifié sans espace de tête grâce à _money_long_signed v5a)
    #   - équité < 0 → afficher "Jusqu'à XX%" (cohérent avec callout
    #     qui devient lui aussi un %)
    if len(prs.slides) > 9:
        s9 = prs.slides[9]
        if rec_nb_log > 0:
            _replace_table_cell_text(
                s9, "Table 2", 5, 1, str(rec_nb_log)
            )
        # Équité scénario 1 (AHP 50 pts)
        if equite_aph50_signed < 0:
            txt_eq1 = _equity_refi_pct(
                equite_aph50_signed, investissement_requis_val
            )
        else:
            txt_eq1 = _money_long_signed(equite_aph50_signed)
        _replace_table_cell_text(s9, "Table 2", 9, 4, txt_eq1)
        # Équité scénario 2 (AHP 100 pts)
        if equite_aph100_signed < 0:
            txt_eq2 = _equity_refi_pct(
                equite_aph100_signed, investissement_requis_val
            )
        else:
            txt_eq2 = _money_long_signed(equite_aph100_signed)
        _replace_table_cell_text(s9, "Table 2", 9, 7, txt_eq2)

    # Slide 2 — nb_etages (legacy : déjà fait dans la section slide 3
    # ci-dessus — on garde l'appel pour rétrocompat)
    nb_etages_input = (
        strat.nb_etages if strat.nb_etages and strat.nb_etages > 0 else None
    )
    if nb_etages_input is not None and len(prs.slides) > 2:
        _replace_table_cell_text(
            prs.slides[2], "Table 12", 1, 1, str(nb_etages_input)
        )

    # Slide 11 — Tendances : chart loyer moyen avant/après (Zipplex)
    if len(prs.slides) > 11:
        s11 = prs.slides[11]
        # Récupère loyer projeté depuis va
        va_for_chart = _calc_value_add(rec, strat)
        loyer_actuel_chart = float(va_for_chart.get("loyer_moyen_actuel", 0) or 0)
        loyer_projete_chart = float(
            va_for_chart.get("nouveau_loyer_moyen", 0) or 0
        )
        tendances_secteur_chart = (rec.city or "").strip() or "le secteur"
        zip_moy_chart, _zc = _zipplex_lookup(tendances_secteur_chart)
        # Override manuel possible
        if strat.tendances_moyenne_actuelle > 0:
            zip_moy_chart = strat.tendances_moyenne_actuelle
        if loyer_actuel_chart > 0 or loyer_projete_chart > 0:
            year0 = loyer_actuel_chart if loyer_actuel_chart > 0 else zip_moy_chart
            year2 = (
                loyer_projete_chart
                if loyer_projete_chart > 0
                else zip_moy_chart * 1.5
            )
            _replace_chart_data(
                s11,
                "Chart 20",
                ["Année 0", "Année 2"],
                [("PDM", [float(year0), float(year2)])],
            )

    # Slide 12 — Valeur ajoutée investisseur : chart trajectoire
    # (Capital investi/récupéré + Équité + Valeurs immeuble)
    if len(prs.slides) > 12:
        s12 = prs.slides[12]
        # Récupère valeurs à partir des chiffres déjà calculés
        results_for_chart = (
            json.loads(rec.analysis_results_json)
            if rec.analysis_results_json
            else None
        )
        va_chart = _calc_value_add(rec, strat)
        invest_req_chart = float(
            rec.mdf_preteur_b
            or (results_for_chart or {}).get("mdf_preteur_b", 0)
            or 0
        )
        nouvelle_vm_12, _src = _resolve_valeur_marchande(
            rec, results_for_chart, va_chart
        )
        ancienne_vm_12 = float(va_chart.get("ancienne_valeur_marchande", 0) or 0)
        # Année 0: -invest, Année 2: +invest récupéré (refi), Année 10:
        #   capital récupéré * coef croissance
        if invest_req_chart > 0:
            # Année 0/2/10 — Capital investi/récupéré
            capital_cat = [
                -invest_req_chart,
                invest_req_chart,
                invest_req_chart * 1.5,
            ]
            equite_an2 = float(va_chart.get("delta_valeur", 0) or 0)
            equite_cat = [
                invest_req_chart * 0.5,
                max(0, equite_an2),
                max(0, equite_an2 * 1.2),
            ]
            valeurs_immeuble_cat = [
                ancienne_vm_12 if ancienne_vm_12 > 0 else float(rec.asking_price or 0),
                nouvelle_vm_12 if nouvelle_vm_12 > 0 else float(rec.asking_price or 0) * 1.5,
                (nouvelle_vm_12 if nouvelle_vm_12 > 0 else float(rec.asking_price or 0) * 1.5) * 1.3,
            ]
            _replace_chart_data(
                s12,
                "Chart 6",
                ["Année 0", "Année 2", "Année 10"],
                [
                    ("Capital investi/recupéré", capital_cat),
                    ("Équité dans l'immeuble", equite_cat),
                    ("Valeurs de l'immeubles", valeurs_immeuble_cat),
                ],
            )

    # Apply photo substitutions
    try:
        with open(mapping_path, "r", encoding="utf-8") as f:
            mapping = json.load(f)
        photo_slots = (mapping.get("photos") or {}).get("slots") or []
    except Exception:
        photo_slots = []
    for i, slot in enumerate(photo_slots):
        if i >= len(resolved_photos):
            break
        slide_idx = int(slot.get("slide", 0))
        shape_name = slot.get("shape_name", "")
        if slide_idx >= len(prs.slides):
            continue
        _replace_photo(prs.slides[slide_idx], shape_name, resolved_photos[i])

    # Serialize
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue(), template_version
